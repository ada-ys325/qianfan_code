import json
import shutil
import sys
import tempfile
import unittest
from datetime import datetime
from pathlib import Path


ROOT = Path(__file__).resolve().parents[1]
if str(ROOT) not in sys.path:
    sys.path.insert(0, str(ROOT))

import openpyxl
from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.shared import Pt, RGBColor
from icalendar import Calendar, Event
from openpyxl.styles import Alignment, Font, PatternFill
from pptx import Presentation
from pptx.dml.color import RGBColor as PptRGBColor
from pptx.util import Inches, Pt as PptPt

from dumatebench.evaluator import (
    evaluate,
    evaluate_calendar_no_overlap,
    evaluate_contain,
    evaluate_diff_contain_text,
    evaluate_directory_structure,
    evaluate_docx_font_style,
    evaluate_docx_paragraph_format,
    evaluate_exact_match,
    evaluate_excel_cell_comparator,
    evaluate_excel_cell_number_tolerance,
    evaluate_excel_cell_style,
    evaluate_excel_cell_value,
    evaluate_excel_formula,
    evaluate_excel_sheet_exists,
    evaluate_file_exist,
    evaluate_file_format_valid,
    evaluate_file_not_exist,
    evaluate_files_unchanged,
    evaluate_log_budget,
    evaluate_no_extra_files,
    evaluate_no_unexpected_diff,
    evaluate_not_contain,
    evaluate_pptx_text_style,
)


class EvaluateAllFunctionsTest(unittest.TestCase):
    def setUp(self):
        self.tmp = tempfile.TemporaryDirectory()
        self.root = Path(self.tmp.name)
        self._create_text_fixtures()
        self._create_excel_fixtures()
        self._create_docx_fixture()
        self._create_pptx_fixture()
        self._create_calendar_fixture()
        self._create_email_fixture()

    def tearDown(self):
        self.tmp.cleanup()

    def _create_text_fixtures(self):
        (self.root / "answer.txt").write_text("Project Alpha total is 1,234.\n")
        (self.root / "expected.txt").write_text("Project Alpha total is 1,234.\n")
        (self.root / "before.txt").write_text("Name: Alice\nStatus: Draft\n")
        (self.root / "after.txt").write_text("Name: Alice\nStatus: Final\n")
        (self.root / "data.json").write_text(json.dumps({"ok": True}))
        (self.root / "data.csv").write_text("name,value\nProject Alpha,1\n")
        (self.root / "brief.pdf").write_bytes(b"%PDF-1.4\n%minimal\n")
        (self.root / "run.json").write_text(json.dumps({"total_tokens": 1200, "elapsed_seconds": 33.5}))
        (self.root / "run.log").write_text("total_tokens: 88\nelapsed_seconds: 9.5\n")

        deliverables = self.root / "deliverables"
        deliverables.mkdir()
        (deliverables / "report.txt").write_text("report")
        (deliverables / "notes.tmp").write_text("ignore")

    def _create_excel_fixtures(self):
        workbook = openpyxl.Workbook()
        sheet = workbook.active
        sheet.title = "Summary"
        sheet["A1"] = 100
        sheet["A2"] = 105
        sheet["A3"] = "=SUM(A1:A2)"
        styled = sheet["B2"]
        styled.value = 42
        styled.font = Font(name="Arial", size=14, bold=True, italic=False, color="FF112233")
        styled.fill = PatternFill(fill_type="solid", fgColor="FFCCEEFF")
        styled.alignment = Alignment(horizontal="center", vertical="center")
        styled.number_format = "0.00"
        workbook.create_sheet("RawData")
        workbook.save(self.root / "metrics.xlsx")
        workbook.close()
        shutil.copyfile(self.root / "metrics.xlsx", self.root / "metrics_expected.xlsx")

    def _create_docx_fixture(self):
        document = Document()
        paragraph = document.add_paragraph()
        paragraph.style = document.styles["Normal"]
        paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
        paragraph.paragraph_format.space_before = Pt(6)
        paragraph.paragraph_format.space_after = Pt(12)
        paragraph.paragraph_format.left_indent = Pt(18)
        paragraph.paragraph_format.line_spacing = 1.5
        run = paragraph.add_run("Styled heading Project Alpha")
        run.font.name = "Arial"
        run.font.size = Pt(16)
        run.font.bold = True
        run.font.italic = False
        run.font.underline = False
        run.font.color.rgb = RGBColor(0x11, 0x22, 0x33)
        document.save(self.root / "styled.docx")

    def _create_pptx_fixture(self):
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
        text_frame = box.text_frame
        text_frame.clear()
        paragraph = text_frame.paragraphs[0]
        run = paragraph.add_run()
        run.text = "Styled slide Project Alpha"
        run.font.name = "Arial"
        run.font.size = PptPt(18)
        run.font.bold = True
        run.font.italic = False
        run.font.color.rgb = PptRGBColor(0x11, 0x22, 0x33)
        presentation.save(self.root / "styled.pptx")

    def _create_calendar_fixture(self):
        calendar_dir = self.root / "calendar"
        calendar_dir.mkdir()
        calendar = Calendar()
        calendar.add("prodid", "-//DuMateBench evaluator test//")
        calendar.add("version", "2.0")
        for summary, start, end in [
            ("Morning", datetime(2026, 1, 1, 9, 0), datetime(2026, 1, 1, 10, 0)),
            ("Planning", datetime(2026, 1, 1, 10, 0), datetime(2026, 1, 1, 11, 0)),
        ]:
            event = Event()
            event.add("summary", summary)
            event.add("dtstart", start)
            event.add("dtend", end)
            calendar.add_component(event)
        (calendar_dir / "Alice.ics").write_bytes(calendar.to_ical())
        shutil.copyfile(calendar_dir / "Alice.ics", self.root / "calendar.ics")

    def _create_email_fixture(self):
        account_dir = self.root / "emails" / "Alice"
        account_dir.mkdir(parents=True)
        (account_dir / "message.eml").write_text("Subject: Update\n\nProject Alpha is complete.")

    def test_file_existence_checks(self):
        self.assertTrue(evaluate_file_exist(self.root, {"file": "answer.txt"}))
        self.assertTrue(evaluate_file_not_exist(self.root, {"file": "missing.txt"}))

    def test_contain_and_not_contain_checks(self):
        self.assertTrue(
            evaluate_contain(
                self.root,
                {"doc_type": "txt", "file": "answer.txt", "keywords": ["Project Alpha", "1234"]},
            )
        )
        self.assertTrue(
            evaluate_contain(
                self.root,
                {"doc_type": "email", "username": "Alice", "keywords": ["Project Alpha"]},
            )
        )
        self.assertTrue(
            evaluate_contain(
                self.root,
                {"doc_type": "docx", "file": "styled.docx", "keywords": ["Styled heading"]},
            )
        )
        self.assertTrue(
            evaluate_contain(
                self.root,
                {"doc_type": "json", "file": "data.json", "keywords": ["ok"]},
            )
        )
        self.assertTrue(
            evaluate_contain(
                self.root,
                {"doc_type": "csv", "file": "data.csv", "keywords": ["Project Alpha"]},
            )
        )
        self.assertTrue(
            evaluate_contain(
                self.root,
                {"doc_type": "pptx", "file": "styled.pptx", "keywords": ["Styled slide"]},
            )
        )
        self.assertTrue(
            evaluate_not_contain(
                self.root,
                {"doc_type": "txt", "file": "answer.txt", "keywords": ["Beta"]},
            )
        )

    def test_file_format_valid_checks_common_supported_types(self):
        for file_name in [
            "answer.txt",
            "data.json",
            "brief.pdf",
            "metrics.xlsx",
            "styled.docx",
            "styled.pptx",
            "calendar.ics",
        ]:
            self.assertTrue(evaluate_file_format_valid(self.root, {"file": file_name}), file_name)

    def test_exact_match_and_diff_checks(self):
        self.assertTrue(
            evaluate_exact_match(
                self.root,
                {"result_file": "answer.txt", "expected_file": "expected.txt", "doc_type": "txt"},
            )
        )
        self.assertTrue(
            evaluate_exact_match(
                self.root,
                {"result_file": "metrics.xlsx", "expected_file": "metrics_expected.xlsx", "doc_type": "xlsx"},
            )
        )
        self.assertTrue(
            evaluate_diff_contain_text(
                self.root,
                {
                    "doc_type": "txt",
                    "input_file": "before.txt",
                    "output_file": "after.txt",
                    "keywords": ["-Status: Draft", "+Status: Final"],
                },
            )
        )
        self.assertTrue(
            evaluate_no_unexpected_diff(
                self.root,
                {
                    "input_file": "before.txt",
                    "output_file": "after.txt",
                    "allowed_patterns": [r"^Status: "],
                },
            )
        )

    def test_files_and_directory_shape_checks(self):
        self.assertTrue(
            evaluate_files_unchanged(
                self.root,
                {"matches": [{"file": "answer.txt", "reference_file": "expected.txt"}]},
            )
        )
        self.assertTrue(
            evaluate_directory_structure(
                self.root,
                {
                    # ``required_files`` are testbed-relative; ignore root.
                    "root": "/",
                    "required_files": ["deliverables/report.txt"],
                    "required_dirs": [],
                    "forbidden_paths": ["deliverables/old.txt"],
                },
            )
        )
        self.assertTrue(
            evaluate_no_extra_files(
                self.root,
                {
                    "root": "deliverables",
                    "expected_files": ["report.txt"],
                    "ignore_patterns": ["*.tmp"],
                },
            )
        )

    def test_excel_value_formula_sheet_numeric_comparator_and_style_checks(self):
        self.assertTrue(
            evaluate_excel_cell_value(
                self.root,
                {"file": "metrics.xlsx", "matches": [{"row": 1, "col": 1, "value": 100}]},
            )
        )
        self.assertTrue(
            evaluate_excel_cell_number_tolerance(
                self.root,
                {"file": "metrics.xlsx", "matches": [{"row": 2, "col": 1, "value": 105.1, "tolerance": 0.2}]},
            )
        )
        self.assertTrue(
            evaluate_excel_sheet_exists(
                self.root,
                {"file": "metrics.xlsx", "required_sheets": ["Summary", "RawData"], "forbidden_sheets": ["Old"]},
            )
        )
        self.assertTrue(
            evaluate_excel_formula(
                self.root,
                {"file": "metrics.xlsx", "matches": [{"sheet": "Summary", "row": 3, "col": 1, "formula": "=SUM(A1:A2)"}]},
            )
        )
        self.assertTrue(
            evaluate_excel_cell_comparator(
                self.root,
                {"file": "metrics.xlsx", "matches": [{"row": 2, "col": 2, "comparator": "lambda x: float(x) > 40"}]},
            )
        )
        self.assertTrue(
            evaluate_excel_cell_style(
                self.root,
                {
                    "file": "metrics.xlsx",
                    "matches": [
                        {
                            "sheet": "Summary",
                            "row": 2,
                            "col": 2,
                            "font_name": "Arial",
                            "font_size_pt": 14,
                            "bold": True,
                            "italic": False,
                            "font_color_rgb": "112233",
                            "fill_color_rgb": "CCEEFF",
                            "number_format": "0.00",
                            "horizontal_alignment": "center",
                            "vertical_alignment": "center",
                        }
                    ],
                },
            )
        )

    def test_docx_style_checks(self):
        self.assertTrue(
            evaluate_docx_font_style(
                self.root,
                {
                    "file": "styled.docx",
                    "matches": [
                        {
                            "paragraph_number": 1,
                            "run_text": "Styled",
                            "font_name": "Arial",
                            "font_size_pt": 16,
                            "bold": True,
                            "italic": False,
                            "underline": False,
                            "color_rgb": "112233",
                        }
                    ],
                },
            )
        )
        self.assertTrue(
            evaluate_docx_paragraph_format(
                self.root,
                {
                    "file": "styled.docx",
                    "matches": [
                        {
                            "text": "Styled heading",
                            "alignment": "CENTER",
                            "left_indent_pt": 18,
                            "space_before_pt": 6,
                            "space_after_pt": 12,
                            "line_spacing": 1.5,
                            "style_name": "Normal",
                        }
                    ],
                },
            )
        )

    def test_pptx_text_style_check(self):
        self.assertTrue(
            evaluate_pptx_text_style(
                self.root,
                {
                    "file": "styled.pptx",
                    "matches": [
                        {
                            "slide_number": 1,
                            "text": "Styled slide",
                            "run_text": "Project Alpha",
                            "font_name": "Arial",
                            "font_size_pt": 18,
                            "bold": True,
                            "italic": False,
                            "color_rgb": "112233",
                        }
                    ],
                },
            )
        )

    def test_log_budget_and_calendar_checks(self):
        self.assertTrue(
            evaluate_log_budget(
                self.root,
                {"log_file": "run.json", "max_tokens": 1500, "max_time_seconds": 40},
            )
        )
        self.assertTrue(
            evaluate_log_budget(
                self.root,
                {"log_file": "run.log", "max_tokens": 100, "max_time_seconds": 10},
            )
        )
        self.assertTrue(evaluate_calendar_no_overlap(self.root, {"username": "Alice"}))

    def test_dispatcher(self):
        self.assertTrue(evaluate("hello Project Alpha", "contain_text", {"keywords": ["project alpha"]}))
        with self.assertRaises(ValueError):
            evaluate("hello", "unsupported", {"keywords": ["hello"]})


if __name__ == "__main__":
    unittest.main(verbosity=2)
