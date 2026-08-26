"""OfficeBench/OdysseyBench-style evaluator functions.

This module copies the evaluator surface used by OfficeBench and OdysseyBench
into DuMateBench while keeping the implementation independent from the
OdysseyBench ``apps.*`` helper modules.
"""

from __future__ import annotations

import difflib
import fnmatch
import glob
import hashlib
import html
import json
import logging
import os
import re
import zipfile
import xml.etree.ElementTree as ET
from pathlib import Path
from typing import Any, Callable

logger = logging.getLogger(__name__)


class UnsupportedArtifactType(ValueError):
    """Raised when a checklist requests a reader that the evaluator does not provide."""


def _is_number(string: Any) -> bool:
    """Check whether a value is numeric (int or float).
    判断一个值是否为数字（整数或浮点数）。"""
    try:
        float(string)
        return True
    except (TypeError, ValueError):
        return False


def _evaluate_contain_text(content: str, args: dict[str, Any]) -> bool:
    """Check whether the text contains all required keywords.
    检查文本内容是否包含所有指定的关键词。"""
    content = content.lower()
    for keyword in args["keywords"]:
        keyword = str(keyword).lower()
        if _is_number(keyword):
            content = content.replace(",", "")
        if keyword not in content:
            return False
    return True


def _read_text(path: str | os.PathLike[str]) -> str:
    """Read a plain text file.
    读取文本文件内容。"""
    return Path(path).read_text(errors="ignore")


def _read_json(path: str | os.PathLike[str]) -> str:
    """Read JSON into a stable, searchable representation."""
    data = json.loads(Path(path).read_text(encoding="utf-8"))
    return json.dumps(data, ensure_ascii=False, sort_keys=True, indent=2)


def _read_html(path: str | os.PathLike[str]) -> str:
    """Extract visible-ish text from HTML without adding a parser dependency."""
    source = Path(path).read_text(errors="ignore")
    source = re.sub(r"(?is)<(script|style|noscript)\b[^>]*>.*?</\1\s*>", " ", source)
    source = re.sub(r"(?s)<[^>]+>", " ", source)
    return re.sub(r"\s+", " ", html.unescape(source)).strip()


def _read_svg(path: str | os.PathLike[str]) -> str:
    """Return SVG text and selected semantic attributes for content checks."""
    try:
        root = ET.parse(path).getroot()
    except ET.ParseError as exc:
        raise ValueError(f"invalid SVG: {path}") from exc
    parts: list[str] = []
    for element in root.iter():
        if element.text and element.text.strip():
            parts.append(element.text.strip())
        for key in ("id", "title", "aria-label", "desc"):
            value = element.attrib.get(key)
            if value:
                parts.append(value)
    return "\n".join(parts)


def _read_zip(path: str | os.PathLike[str]) -> str:
    """Return a ZIP member inventory plus safely decodable text members."""
    parts: list[str] = []
    with zipfile.ZipFile(path) as archive:
        for info in archive.infolist():
            if info.is_dir():
                continue
            parts.append(info.filename)
            if info.file_size > 1_000_000 or Path(info.filename).suffix.lower() not in {
                ".txt", ".md", ".json", ".csv", ".yaml", ".yml", ".py", ".js", ".java",
            }:
                continue
            try:
                parts.append(archive.read(info).decode("utf-8", errors="ignore"))
            except (OSError, zipfile.BadZipFile):
                continue
    return "\n".join(parts)


def _read_xlsx(path: str | os.PathLike[str]) -> str:
    """Read an Excel file and return the position and value of each non-empty cell.
    读取 Excel 文件内容，返回每个非空单元格的位置和值。"""
    try:
        import openpyxl
    except ImportError as exc:
        raise RuntimeError("openpyxl is required to evaluate xlsx files") from exc

    workbook = openpyxl.load_workbook(path, data_only=True)
    parts: list[str] = []
    for sheet in workbook.worksheets:
        for row in sheet.iter_rows():
            for cell in row:
                if cell.value is not None:
                    parts.append(f"({cell.row}, {cell.column}): {cell.value}")
    return "\n".join(parts) + "\n"


def _read_docx(path: str | os.PathLike[str]) -> str:
    """Read a Word document and return the text of all paragraphs.
    读取 Word 文档内容，返回所有段落的文本。"""
    try:
        from docx import Document
    except ImportError as exc:
        raise RuntimeError("python-docx is required to evaluate docx files") from exc

    document = Document(path)
    parts = [paragraph.text for paragraph in document.paragraphs if paragraph.text]
    for table in document.tables:
        for row in table.rows:
            parts.append(" | ".join(cell.text for cell in row.cells))
    for section in document.sections:
        parts.extend(paragraph.text for paragraph in section.header.paragraphs if paragraph.text)
        parts.extend(paragraph.text for paragraph in section.footer.paragraphs if paragraph.text)
    return "\n".join(parts)


def _read_pptx(path: str | os.PathLike[str]) -> str:
    """Extract slide text, table cells and speaker notes from a PowerPoint file."""
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise RuntimeError("python-pptx is required to inspect pptx files") from exc

    presentation = Presentation(path)
    parts: list[str] = []
    for slide_number, slide in enumerate(presentation.slides, start=1):
        parts.append(f"[SLIDE {slide_number}]")
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False) and shape.text:
                parts.append(shape.text)
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    parts.append(" | ".join(cell.text for cell in row.cells))
        try:
            notes = slide.notes_slide.notes_text_frame
        except (AttributeError, KeyError):
            notes = None
        if notes is not None and notes.text:
            parts.append(f"[NOTES] {notes.text}")
    return "\n".join(parts)


def _read_pdf(path: str | os.PathLike[str]) -> str:
    """Read a PDF file and return extracted text from all pages.
    读取 PDF 文件内容，返回所有页面的提取文本。"""
    try:
        from pypdf import PdfReader
    except ImportError:
        try:
            from PyPDF2 import PdfReader  # type: ignore
        except ImportError as exc:
            raise RuntimeError("pypdf or PyPDF2 is required to evaluate pdf files") from exc

    reader = PdfReader(path)
    return "\n".join(page.extract_text() or "" for page in reader.pages)


def _read_email_account(account_dir: str | os.PathLike[str]) -> str:
    """Read all .eml files under an email account directory.
    读取邮件账户目录下所有 .eml 文件的内容。"""
    messages: list[str] = []
    for file_path in sorted(Path(account_dir).glob("*.eml")):
        messages.append(file_path.read_text(errors="ignore"))
    return "\n".join(messages)


def _task_path(testbed_dir: str | os.PathLike[str], file_path: str) -> Path:
    """Join the testbed directory and file path into a Path object.
    拼接测试目录和文件路径，返回完整的 Path 对象。"""
    return Path(testbed_dir) / file_path


def _file_sha256(file_path: str | os.PathLike[str]) -> str:
    """Compute the SHA-256 hash of a file.
    计算文件的 SHA-256 哈希值。"""
    digest = hashlib.sha256()
    with Path(file_path).open("rb") as handle:
        for chunk in iter(lambda: handle.read(1024 * 1024), b""):
            digest.update(chunk)
    return digest.hexdigest()


def _relative_files(root: Path) -> set[str]:
    """Get the set of relative file paths under a directory.
    获取目录下所有文件的相对路径集合。"""
    if not root.exists() or not root.is_dir():
        return set()
    return {
        path.relative_to(root).as_posix()
        for path in root.rglob("*")
        if path.is_file()
    }


def _matches_any(path: str, patterns: list[str]) -> bool:
    """Check whether a path matches any of the glob patterns.
    检查路径是否匹配任意一个 glob 模式。"""
    return any(fnmatch.fnmatch(path, pattern) for pattern in patterns)


def _none_safe_equal(actual: Any, expected: Any) -> bool:
    return actual is not None and actual == expected


def _points_equal(actual: Any, expected: Any, tolerance: float = 0.1) -> bool:
    if actual is None:
        return False
    try:
        actual_pt = actual.pt if hasattr(actual, "pt") else float(actual)
        return abs(float(actual_pt) - float(expected)) <= tolerance
    except (TypeError, ValueError):
        return False


def _normalize_rgb(value: Any) -> str | None:
    if value is None:
        return None
    text = str(value).upper()
    if text in {"NONE", "AUTO"}:
        return None
    return text[-6:]


def _docx_paragraph_matches(paragraph: Any, match: dict[str, Any]) -> bool:
    if "text" not in match:
        return True
    return str(match["text"]) in paragraph.text


def _docx_target_paragraphs(document: Any, match: dict[str, Any]) -> list[Any]:
    paragraphs = list(document.paragraphs)
    if "paragraph_number" in match:
        index = int(match["paragraph_number"]) - 1
        return [paragraphs[index]] if 0 <= index < len(paragraphs) else []
    if "paragraph_index" in match:
        index = int(match["paragraph_index"])
        return [paragraphs[index]] if 0 <= index < len(paragraphs) else []
    return [paragraph for paragraph in paragraphs if _docx_paragraph_matches(paragraph, match)]


def _reader_for_doc_type(doc_type: str) -> Callable[[str | os.PathLike[str]], str]:
    """Return the content reader function for the given document type.
    根据文档类型返回对应的内容读取函数。"""
    doc_type = str(doc_type).lower().lstrip(".")
    if doc_type == "xlsx":
        return _read_xlsx
    if doc_type in {
        "txt", "ics", "md", "csv", "tsv", "yaml", "yml", "py", "js", "ts", "java",
        "xml", "ini", "cfg", "log",
    }:
        return _read_text
    if doc_type == "json":
        return _read_json
    if doc_type in {"html", "htm"}:
        return _read_html
    if doc_type in {"doc", "docx"}:
        return _read_docx
    if doc_type in {"ppt", "pptx"}:
        return _read_pptx
    if doc_type == "pdf":
        return _read_pdf
    if doc_type == "svg":
        return _read_svg
    if doc_type == "zip":
        return _read_zip
    raise UnsupportedArtifactType(f"Unsupported document type for content checks: {doc_type}")


def evaluate_contain(testbed_dir: str | os.PathLike[str], args: dict[str, Any]) -> bool:
    """Check whether a document or email account contains all keywords.
    检查文档或邮件账户是否包含所有关键词。"""
    doc_type = str(args.get("doc_type") or Path(str(args.get("file", ""))).suffix.lstrip("."))
    testbed_dir = str(testbed_dir)

    if doc_type == "email":
        username = args["username"]
        email_contents = ""
        exact_dir = os.path.join(testbed_dir, "emails", username)
        lower_dir = os.path.join(testbed_dir, "emails", username.lower())
        if os.path.exists(exact_dir):
            email_contents = _read_email_account(exact_dir)
        elif os.path.exists(lower_dir):
            email_contents = _read_email_account(lower_dir)
        else:
            for email_account in glob.glob(os.path.join(testbed_dir, "emails", "*")):
                if username.lower() in os.path.basename(email_account).lower():
                    email_contents = _read_email_account(email_account)
                    break
        return _evaluate_contain_text(email_contents, args)

    file_path = os.path.join(testbed_dir, args["file"])
    if not os.path.isfile(file_path):
        logger.debug("File does not exist: %s", file_path)
        return False
    content = _reader_for_doc_type(doc_type)(file_path)
    return _evaluate_contain_text(content, args)


def evaluate_not_contain(testbed_dir: str | os.PathLike[str], args: dict[str, Any]) -> bool:
    """Check whether a document or email account misses at least one keyword.
    检查文档或邮件账户是否缺少至少一个关键词。"""
    if args.get("doc_type") != "email":
        file_path = _task_path(testbed_dir, str(args.get("file", "")))
        if not file_path.is_file():
            logger.debug("File does not exist: %s", file_path)
            return False
    return not evaluate_contain(testbed_dir, args)


def evaluate_file_exist(testbed_dir: str | os.PathLike[str], args: dict[str, Any]) -> bool:
    """Check whether a file or directory exists under the testbed.
    检查测试目录下是否存在指定文件或目录。"""
    return os.path.exists(os.path.join(testbed_dir, args["file"]))


def evaluate_file_not_exist(testbed_dir: str | os.PathLike[str], args: dict[str, Any]) -> bool:
    """Check whether a file or directory does not exist under the testbed.
    检查测试目录下是否不存在指定文件或目录。"""
    return not os.path.exists(os.path.join(testbed_dir, args["file"]))


def evaluate_file_format_valid(
    testbed_dir: str | os.PathLike[str], args: dict[str, Any]
) -> bool:
    """Check whether a file is structurally readable for its declared type.
    检查文件在结构上能否被其声明类型正常读取。"""
    file_path = _task_path(testbed_dir, args["file"])
    if not file_path.is_file():
        return False

    doc_type = str(args.get("doc_type") or file_path.suffix.lower().lstrip(".")).lower().lstrip(".")
    try:
        if doc_type in {"txt", "md", "csv", "tsv", "yaml", "yml", "py", "js", "ts", "java", "html", "htm", "xml"}:
            file_path.read_text(errors="strict")
            return True
        if doc_type == "json":
            json.loads(file_path.read_text())
            return True
        if doc_type in {"docx", "pptx"}:
            marker = "word/document.xml" if doc_type == "docx" else "ppt/presentation.xml"
            with zipfile.ZipFile(file_path) as archive:
                return marker in archive.namelist()
        if doc_type == "xlsx":
            try:
                import openpyxl
            except ImportError as exc:
                raise RuntimeError("openpyxl is required to validate xlsx files") from exc
            openpyxl.load_workbook(file_path, read_only=True, data_only=False).close()
            return True
        if doc_type == "pdf":
            return file_path.read_bytes().startswith(b"%PDF")
        if doc_type == "svg":
            ET.parse(file_path)
            return True
        if doc_type == "zip":
            with zipfile.ZipFile(file_path) as archive:
                return archive.testzip() is None
        if doc_type == "ics":
            content = file_path.read_text(errors="ignore")
            return "BEGIN:VCALENDAR" in content and "END:VCALENDAR" in content
    except RuntimeError:
        raise
    except Exception as exc:
        logger.debug("Format validation failed for %s: %s", file_path, exc)
        return False

    logger.debug("Unsupported doc_type for format validation: %s", doc_type)
    return False


def evaluate_files_unchanged(
    testbed_dir: str | os.PathLike[str], args: dict[str, Any]
) -> bool:
    """Check whether files match their reference copies byte-for-byte.
    检查文件是否与参考副本逐字节一致。"""
    for match in args["matches"]:
        result_file = _task_path(testbed_dir, match["file"])
        reference_file = _task_path(testbed_dir, match["reference_file"])
        if not result_file.is_file() or not reference_file.is_file():
            return False
        if _file_sha256(result_file) != _file_sha256(reference_file):
            return False
    return True


def evaluate_no_unexpected_diff(
    testbed_dir: str | os.PathLike[str], args: dict[str, Any]
) -> bool:
    """Check that changed diff lines are limited to allowed regex patterns.
    检查文件的变更行是否仅限于允许的正则表达式模式。"""
    input_file = _task_path(testbed_dir, args["input_file"])
    output_file = _task_path(testbed_dir, args["output_file"])
    if not input_file.is_file() or not output_file.is_file():
        return False

    input_lines = input_file.read_text(errors="ignore").splitlines()
    output_lines = output_file.read_text(errors="ignore").splitlines()
    allowed_patterns = [re.compile(pattern) for pattern in args.get("allowed_patterns", [])]
    diff = difflib.unified_diff(input_lines, output_lines, n=0)

    for line in diff:
        if not line or line.startswith(("+++", "---", "@@")):
            continue
        if not line.startswith(("+", "-")):
            continue
        changed_text = line[1:]
        if not any(pattern.search(changed_text) for pattern in allowed_patterns):
            return False
    return True


def evaluate_directory_structure(
    testbed_dir: str | os.PathLike[str], args: dict[str, Any]
) -> bool:
    """Check required and forbidden paths directly under the testbed.

    ``required_files`` and the other path lists already contain paths relative
    to the evaluator testbed (commonly ``run_outputs/...``).  The historical
    implementation prepended ``args['root']`` to every entry, which made a
    valid check such as ``root: /`` or ``required_files: [run_outputs/a]``
    resolve outside the task.  Match ``evaluate_file_exist`` semantics and
    resolve each listed path directly against ``testbed_dir``; ``root`` is
    intentionally ignored for backwards-compatible checklist arguments.
    """
    for file_name in args.get("required_files", []):
        if not (_task_path(testbed_dir, file_name)).is_file():
            return False
    for dir_name in args.get("required_dirs", []):
        if not (_task_path(testbed_dir, dir_name)).is_dir():
            return False
    for path_name in args.get("forbidden_paths", []):
        if _task_path(testbed_dir, path_name).exists():
            return False
    return True


def evaluate_no_extra_files(
    testbed_dir: str | os.PathLike[str], args: dict[str, Any]
) -> bool:
    """Check that a directory contains no files outside an expected set.
    检查目录中是否不包含预期集合之外的多余文件。"""
    root = _task_path(testbed_dir, args.get("root", "."))
    expected = {Path(path).as_posix() for path in args.get("expected_files", [])}
    ignore_patterns = args.get("ignore_patterns", [])
    actual = {
        path for path in _relative_files(root)
        if not _matches_any(path, ignore_patterns)
    }
    return actual.issubset(expected)


def _helper_diff_contain_text(
    input_content: str, output_content: str, args: dict[str, Any]
) -> bool:
    """Check whether the unified diff of two texts contains all keywords.
    检查两个文本的 unified diff 是否包含所有关键词。"""
    if input_content == output_content:
        return False
    diff = difflib.unified_diff(input_content.split("\n"), output_content.split("\n"), n=0)
    diff_text = "\n".join(list(diff))
    return all(match in diff_text for match in args["keywords"])


def evaluate_diff_contain_text(
    testbed_dir: str | os.PathLike[str], args: dict[str, Any]
) -> bool:
    """Check whether the unified diff of two files contains all keywords.
    检查两个文件的 unified diff 是否包含所有关键词。"""
    doc_type = args["doc_type"]
    input_file = os.path.join(testbed_dir, args["input_file"])
    output_file = os.path.join(testbed_dir, args["output_file"])
    helper = _reader_for_doc_type(doc_type)
    return _helper_diff_contain_text(helper(input_file), helper(output_file), args)


def evaluate_excel_cell_value(
    testbed_dir: str | os.PathLike[str], args: dict[str, Any]
) -> bool:
    """Check whether specified Excel cells match exact stringified values.
    检查指定的 Excel 单元格是否与预期的字符串化值完全匹配。"""
    file_name = args.get("file") or args.get("output_file")
    file_path = os.path.join(testbed_dir, file_name)
    if not os.path.exists(file_path):
        logger.debug("File does not exist: %s", file_path)
        return False
    content = _read_xlsx(file_path)
    for match in args["matches"]:
        pattern = f'({match["row"]}, {match["col"]}): {match["value"]}'
        if pattern not in content:
            return False
    return True


def evaluate_excel_cell_number_tolerance(
    testbed_dir: str | os.PathLike[str], args: dict[str, Any]
) -> bool:
    """Check numeric Excel cells within absolute tolerance.
    检查 Excel 数值单元格是否在指定的绝对容差范围内。"""
    try:
        import openpyxl
    except ImportError:
        logger.warning("openpyxl is required to evaluate xlsx files")
        return False

    file_name = args.get("file") or args.get("output_file")
    file_path = _task_path(testbed_dir, file_name)
    if not file_path.exists():
        return False
    workbook = openpyxl.load_workbook(file_path, data_only=True)
    try:
        for match in args["matches"]:
            sheet = workbook[match["sheet"]] if match.get("sheet") else workbook.active
            value = sheet.cell(row=match["row"], column=match["col"]).value
            expected = float(match["value"])
            tolerance = float(match.get("tolerance", args.get("tolerance", 0)))
            if value is None or abs(float(value) - expected) > tolerance:
                return False
    finally:
        workbook.close()
    return True


def evaluate_excel_sheet_exists(
    testbed_dir: str | os.PathLike[str], args: dict[str, Any]
) -> bool:
    """Check required and forbidden sheet names in an Excel workbook.
    检查 Excel 工作簿中必需的（和禁止的）工作表名称。"""
    try:
        import openpyxl
    except ImportError:
        logger.warning("openpyxl is required to evaluate xlsx files")
        return False

    file_path = _task_path(testbed_dir, args["file"])
    if not file_path.exists():
        return False
    workbook = openpyxl.load_workbook(file_path, read_only=True, data_only=False)
    try:
        sheet_names = set(workbook.sheetnames)
        required = set(args.get("required_sheets", []))
        forbidden = set(args.get("forbidden_sheets", []))
        return required.issubset(sheet_names) and sheet_names.isdisjoint(forbidden)
    finally:
        workbook.close()


def evaluate_excel_formula(
    testbed_dir: str | os.PathLike[str], args: dict[str, Any]
) -> bool:
    """Check exact formula strings in Excel cells.
    检查 Excel 单元格中的公式字符串是否精确匹配。"""
    try:
        import openpyxl
    except ImportError:
        logger.warning("openpyxl is required to evaluate xlsx files")
        return False

    file_path = _task_path(testbed_dir, args["file"])
    if not file_path.exists():
        return False
    workbook = openpyxl.load_workbook(file_path, data_only=False)
    try:
        for match in args["matches"]:
            sheet = workbook[match["sheet"]] if match.get("sheet") else workbook.active
            value = sheet.cell(row=match["row"], column=match["col"]).value
            if value != match["formula"]:
                return False
    finally:
        workbook.close()
    return True


def evaluate_docx_font_style(
    testbed_dir: str | os.PathLike[str], args: dict[str, Any]
) -> bool:
    """Check font settings in Word document runs.
    检查 Word 文档 run 级别的字体设置。"""
    try:
        from docx import Document
    except ImportError:
        logger.warning("python-docx is required to evaluate docx font style")
        return False

    file_path = _task_path(testbed_dir, args["file"])
    if not file_path.exists():
        return False
    document = Document(file_path)

    for match in args["matches"]:
        paragraphs = _docx_target_paragraphs(document, match)
        if not paragraphs:
            return False
        passed = False
        for paragraph in paragraphs:
            runs = list(paragraph.runs)
            if "run_index" in match:
                index = int(match["run_index"])
                runs = [runs[index]] if 0 <= index < len(runs) else []
            elif "run_text" in match:
                runs = [run for run in runs if str(match["run_text"]) in run.text]
            for run in runs:
                font = run.font
                checks = [
                    "font_name" not in match or _none_safe_equal(font.name, match["font_name"]),
                    "font_size_pt" not in match or _points_equal(font.size, match["font_size_pt"]),
                    "bold" not in match or font.bold is bool(match["bold"]),
                    "italic" not in match or font.italic is bool(match["italic"]),
                    "underline" not in match or font.underline is bool(match["underline"]),
                    "color_rgb" not in match
                    or _normalize_rgb(font.color.rgb) == _normalize_rgb(match["color_rgb"]),
                ]
                if all(checks):
                    passed = True
                    break
            if passed:
                break
        if not passed:
            return False
    return True


def evaluate_docx_paragraph_format(
    testbed_dir: str | os.PathLike[str], args: dict[str, Any]
) -> bool:
    """Check paragraph-level formatting in a Word document.
    检查 Word 文档段落级别的格式设置。"""
    try:
        from docx import Document
    except ImportError:
        logger.warning("python-docx is required to evaluate docx paragraph format")
        return False

    file_path = _task_path(testbed_dir, args["file"])
    if not file_path.exists():
        return False
    document = Document(file_path)

    for match in args["matches"]:
        paragraphs = _docx_target_paragraphs(document, match)
        if not paragraphs:
            return False
        passed = False
        for paragraph in paragraphs:
            fmt = paragraph.paragraph_format
            alignment = paragraph.alignment.name if paragraph.alignment is not None else None
            checks = [
                "style_name" not in match or paragraph.style.name == match["style_name"],
                "alignment" not in match or alignment == str(match["alignment"]).upper(),
                "left_indent_pt" not in match
                or _points_equal(fmt.left_indent, match["left_indent_pt"]),
                "right_indent_pt" not in match
                or _points_equal(fmt.right_indent, match["right_indent_pt"]),
                "first_line_indent_pt" not in match
                or _points_equal(fmt.first_line_indent, match["first_line_indent_pt"]),
                "space_before_pt" not in match
                or _points_equal(fmt.space_before, match["space_before_pt"]),
                "space_after_pt" not in match
                or _points_equal(fmt.space_after, match["space_after_pt"]),
                "line_spacing" not in match
                or fmt.line_spacing is not None
                and abs(float(fmt.line_spacing) - float(match["line_spacing"])) <= 0.01,
            ]
            if all(checks):
                passed = True
                break
        if not passed:
            return False
    return True


def evaluate_excel_cell_style(
    testbed_dir: str | os.PathLike[str], args: dict[str, Any]
) -> bool:
    """Check Excel cell font, fill, alignment, and number format.
    检查 Excel 单元格的字体、填充、对齐和数字格式。"""
    try:
        import openpyxl
    except ImportError:
        logger.warning("openpyxl is required to evaluate xlsx styles")
        return False

    file_path = _task_path(testbed_dir, args["file"])
    if not file_path.exists():
        return False
    workbook = openpyxl.load_workbook(file_path, data_only=False)
    try:
        for match in args["matches"]:
            sheet = workbook[match["sheet"]] if match.get("sheet") else workbook.active
            cell = sheet.cell(row=match["row"], column=match["col"])
            checks = [
                "font_name" not in match or cell.font.name == match["font_name"],
                "font_size_pt" not in match
                or abs(float(cell.font.sz) - float(match["font_size_pt"])) <= 0.1,
                "bold" not in match or cell.font.bold is bool(match["bold"]),
                "italic" not in match or cell.font.italic is bool(match["italic"]),
                "underline" not in match or bool(cell.font.underline) is bool(match["underline"]),
                "font_color_rgb" not in match
                or _normalize_rgb(cell.font.color.rgb if cell.font.color else None)
                == _normalize_rgb(match["font_color_rgb"]),
                "fill_color_rgb" not in match
                or _normalize_rgb(cell.fill.fgColor.rgb) == _normalize_rgb(match["fill_color_rgb"]),
                "number_format" not in match or cell.number_format == match["number_format"],
                "horizontal_alignment" not in match
                or cell.alignment.horizontal == match["horizontal_alignment"],
                "vertical_alignment" not in match
                or cell.alignment.vertical == match["vertical_alignment"],
            ]
            if not all(checks):
                return False
    finally:
        workbook.close()
    return True


def evaluate_pptx_text_style(
    testbed_dir: str | os.PathLike[str], args: dict[str, Any]
) -> bool:
    """Check text run font settings in PowerPoint slides.
    检查 PowerPoint 幻灯片文本 run 的字体设置。"""
    try:
        from pptx import Presentation
    except ImportError:
        logger.warning("python-pptx is required to evaluate pptx text style")
        return False

    file_path = _task_path(testbed_dir, args["file"])
    if not file_path.exists():
        return False
    presentation = Presentation(file_path)

    for match in args["matches"]:
        if "slide_number" in match:
            slide_indexes = [int(match["slide_number"]) - 1]
        elif "slide_index" in match:
            slide_indexes = [int(match["slide_index"])]
        else:
            slide_indexes = list(range(len(presentation.slides)))

        passed = False
        for slide_index in slide_indexes:
            if not 0 <= slide_index < len(presentation.slides):
                continue
            slide = presentation.slides[slide_index]
            for shape in slide.shapes:
                if not getattr(shape, "has_text_frame", False):
                    continue
                if "text" in match and str(match["text"]) not in shape.text:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if "run_text" in match and str(match["run_text"]) not in run.text:
                            continue
                        font = run.font
                        checks = [
                            "font_name" not in match or font.name == match["font_name"],
                            "font_size_pt" not in match or _points_equal(font.size, match["font_size_pt"]),
                            "bold" not in match or font.bold is bool(match["bold"]),
                            "italic" not in match or font.italic is bool(match["italic"]),
                            "color_rgb" not in match
                            or _normalize_rgb(font.color.rgb) == _normalize_rgb(match["color_rgb"]),
                        ]
                        if all(checks):
                            passed = True
                            break
                    if passed:
                        break
                if passed:
                    break
            if passed:
                break
        if not passed:
            return False
    return True


def evaluate_pptx_llm_judge(
    testbed_dir: str | os.PathLike[str], args: dict[str, Any]
) -> bool:
    """Evaluate a PowerPoint artifact with the PPT LLM-as-judge.
    使用 PPT 方向的 LLM-as-judge 评估 PowerPoint 产物质量。"""
    try:
        from .llm_judge.ppt import evaluate_pptx_llm_judge as _evaluate
    except ImportError:
        try:
            from dumatebench.evaluator.llm_judge.ppt import evaluate_pptx_llm_judge as _evaluate
        except ImportError as exc:
            logger.warning("PPT LLM judge import failed: %s", exc)
            return False
    return _evaluate(testbed_dir, args)


def _check_eval_name(check: dict[str, Any]) -> str:
    name = str(check.get("function") or check.get("type") or "")
    aliases = {
        "file_exists": "evaluate_file_exist",
        "file_not_exists": "evaluate_file_not_exist",
        "file_not_exist": "evaluate_file_not_exist",
        "file_format_valid": "evaluate_file_format_valid",
    }
    return aliases.get(name, name)


def run_checklist_score(
    testbed_dir: str | os.PathLike[str],
    checks: list[dict[str, Any]],
) -> dict[str, Any]:
    """Run objective checks while preserving evaluator errors distinct from failures.

    Existing callers can keep consuming ``passed``. New callers should use
    ``status`` and ``score_eligible`` so an unsupported reader, missing
    dependency, or malformed check is never reported as an Agent failure.
    """
    results: list[dict[str, Any]] = []
    for index, check in enumerate(checks, start=1):
        check_id = str(check.get("id") or f"check_{index}")
        eval_name = _check_eval_name(check)
        eval_func = globals().get(eval_name)
        if not callable(eval_func) or eval_name == "evaluate_llm_judge_score":
            passed = False
            detail = f"unsupported checklist function: {eval_name}"
            status = "unsupported"
            score_eligible = False
        else:
            check_dir = _task_path(testbed_dir, check.get("testbed_dir", "."))
            check_args = check.get("args", {})
            if "path" in check and "file" not in check_args:
                check_args = {**check_args, "file": check["path"]}
            try:
                passed = bool(eval_func(check_dir, check_args))
                detail = eval_name
                status = "pass" if passed else "fail"
                score_eligible = True
            except UnsupportedArtifactType as exc:
                logger.warning("Checklist check %s is unsupported: %s", check_id, exc)
                passed = False
                detail = f"{eval_name}: {type(exc).__name__}: {exc}"
                status = "unsupported"
                score_eligible = False
            except Exception as exc:
                logger.warning("Checklist check %s failed: %s", check_id, exc)
                passed = False
                detail = f"{eval_name}: {type(exc).__name__}: {exc}"
                status = "evaluator_error"
                score_eligible = False
        results.append(
            {
                "id": check_id,
                "type": eval_name,
                "weight": 1.0,
                "passed": passed,
                "detail": detail,
                "status": status,
                "score_eligible": score_eligible,
            }
        )
    eligible = [item for item in results if item["score_eligible"]]
    partial_pass = round(sum(item["passed"] for item in eligible) / len(eligible), 4) if eligible else None
    return {
        "complete_pass": int(bool(results) and len(eligible) == len(results) and all(item["passed"] for item in eligible)),
        "partial_pass": partial_pass,
        "eligible_check_count": len(eligible),
        "ineligible_check_count": len(results) - len(eligible),
        "checks": results,
    }


def run_llm_judge_score(
    testbed_dir: str | os.PathLike[str], args: dict[str, Any]
) -> dict[str, Any]:
    """Run checklist + type-specific LLM judge and return normalized score details."""
    try:
        from .llm_judge.unified import run_llm_judge_score as _run
    except ImportError:
        try:
            from dumatebench.evaluator.llm_judge.unified import run_llm_judge_score as _run
        except ImportError as exc:
            logger.warning("Unified LLM judge import failed: %s", exc)
            raise
    return _run(testbed_dir, args, checklist_runner=run_checklist_score)


def evaluate_llm_judge_score(
    testbed_dir: str | os.PathLike[str], args: dict[str, Any]
) -> bool:
    """Boolean evaluator wrapper for unified checklist + LLM judge scoring."""
    try:
        return bool(run_llm_judge_score(testbed_dir, args).get("pass"))
    except Exception as exc:
        logger.warning("Unified LLM judge failed: %s", exc)
        return False


def evaluate_excel_cell_comparator(
    testbed_dir: str | os.PathLike[str], args: dict[str, Any]
) -> bool:
    """Check whether specified Excel cells satisfy comparator expressions.
    检查指定的 Excel 单元格是否满足比较器表达式。"""
    file_name = args.get("file") or args.get("output_file")
    file_path = os.path.join(testbed_dir, file_name)
    if not os.path.exists(file_path):
        logger.debug("File does not exist: %s", file_path)
        return False
    content = _read_xlsx(file_path)
    for match in args["matches"]:
        pattern = r"\({}, {}\): ([^\t\n]+)[\t\n]".format(match["row"], match["col"])
        value_match = re.search(pattern, content)
        if not value_match:
            return False
        if not eval(match["comparator"])(value_match.group(1)):
            return False
    return True


def evaluate_log_budget(testbed_dir: str | os.PathLike[str], args: dict[str, Any]) -> bool:
    """Check token/time budgets from a JSON log or text log.
    从 JSON 或文本日志中检查 token/时间预算是否超出限制。"""
    log_path = _task_path(testbed_dir, args["log_file"])
    if not log_path.is_file():
        return False
    text = log_path.read_text(errors="ignore")

    values: dict[str, float] = {}
    try:
        data = json.loads(text)
        if isinstance(data, dict):
            for key in ("tokens", "total_tokens", "time_seconds", "elapsed_seconds"):
                if key in data and data[key] is not None:
                    values[key] = float(data[key])
    except json.JSONDecodeError:
        token_match = re.search(r"(?:total_)?tokens?\D+(\d+(?:\.\d+)?)", text, re.I)
        time_match = re.search(r"(?:time|elapsed)(?:_seconds)?\D+(\d+(?:\.\d+)?)", text, re.I)
        if token_match:
            values["tokens"] = float(token_match.group(1))
        if time_match:
            values["time_seconds"] = float(time_match.group(1))

    token_value = values.get("total_tokens", values.get("tokens"))
    time_value = values.get("elapsed_seconds", values.get("time_seconds"))
    if "max_tokens" in args and (token_value is None or token_value > float(args["max_tokens"])):
        return False
    if "max_time_seconds" in args and (
        time_value is None or time_value > float(args["max_time_seconds"])
    ):
        return False
    return True


def evaluate_calendar_no_overlap(
    testbed_dir: str | os.PathLike[str], args: dict[str, Any]
) -> bool:
    """Check whether a user's .ics calendar has no overlapping events.
    检查用户的 .ics 日历中是否存在时间重叠的事件。"""
    try:
        import icalendar
        import pytz
    except ImportError:
        logger.warning("icalendar and pytz are required to evaluate calendars")
        return False

    username = args["username"]
    calendar_file = os.path.join(testbed_dir, "calendar", f"{username}.ics")
    if not os.path.exists(calendar_file):
        logger.debug("Calendar does not exist: %s", calendar_file)
        return False

    calendar = icalendar.Calendar.from_ical(Path(calendar_file).read_bytes())
    utc = pytz.UTC

    def proc_dt(value):
        if value.tzinfo is None:
            return utc.localize(value)
        return value

    events = [component for component in calendar.walk() if component.name == "VEVENT"]
    events.sort(key=lambda event: proc_dt(event.get("dtstart").dt))
    for index in range(len(events) - 1):
        if proc_dt(events[index].get("dtend").dt) > proc_dt(events[index + 1].get("dtstart").dt):
            logger.debug(
                "Calendar of %s: Event %s and Event %s overlap",
                username,
                events[index].get("summary"),
                events[index + 1].get("summary"),
            )
            return False
    return True


def evaluate_exact_match(testbed_dir: str | os.PathLike[str], args: dict[str, Any]) -> bool:
    """Check whether result and expected files match exactly.
    检查结果文件与预期文件是否完全匹配。

    对于非 xlsx 文件，使用对应类型的读取器将两个文件内容读取后直接比较。
    对于 xlsx 文件，逐单元格比较两个工作表的内容（双向校验，确保行列数量一致）。

    Args:
        testbed_dir: 测试目录路径，结果文件和预期文件均位于该目录下。
        args: 评估参数字典，需包含以下键：
            - result_file (str): 结果文件的相对路径。
            - expected_file (str): 预期文件的相对路径。
            - doc_type (str): 文档类型，如 "xlsx"、"docx"、"pdf" 等。

    Returns:
        bool: 文件内容完全匹配时返回 True，否则返回 False。
              若结果文件或预期文件不存在，或 xlsx 依赖缺失，也返回 False。
    """
    result_path = os.path.join(testbed_dir, args["result_file"])
    expected_path = os.path.join(testbed_dir, args["expected_file"])
    if not os.path.exists(result_path) or not os.path.exists(expected_path):
        logger.debug("Missing result or expected file: %s, %s", result_path, expected_path)
        return False

    doc_type = args["doc_type"]
    if doc_type != "xlsx":
        helper = _reader_for_doc_type(doc_type)
        return helper(result_path) == helper(expected_path)

    try:
        import openpyxl
    except ImportError:
        logger.warning("openpyxl is required to evaluate xlsx files")
        return False

    result_sheet = openpyxl.load_workbook(result_path, data_only=True).active
    expected_sheet = openpyxl.load_workbook(expected_path, data_only=True).active

    for row in result_sheet.iter_rows():
        for cell in row:
            expected_value = expected_sheet.cell(row=cell.row, column=cell.column).value
            if cell.value != expected_value:
                return False

    for row in expected_sheet.iter_rows():
        for cell in row:
            result_value = result_sheet.cell(row=cell.row, column=cell.column).value
            if result_value != cell.value:
                return False

    return True


def evaluate(
    testbed_dir: str | os.PathLike[str], evaluate_type: str, args: dict[str, Any]
) -> bool:
    """OfficeBench-compatible dispatcher for primitive evaluator types.
    OfficeBench 兼容的评估器分发函数，根据评估类型调用对应的评估逻辑。"""
    if evaluate_type == "contain_text":
        return _evaluate_contain_text(str(testbed_dir), args)
    raise ValueError(f"Invalid evaluate type: {evaluate_type}")
