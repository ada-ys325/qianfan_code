"""LLM-as-judge evaluator for PowerPoint artifacts."""

from __future__ import annotations

import argparse
import base64
import json
import logging
import os
import shutil
import subprocess
import tempfile
import time
import zipfile
from pathlib import Path
from typing import Any

from ..scoring import equal_weight_partial_pass, final_score as merge_final_score

logger = logging.getLogger(__name__)

DEFAULT_MODEL = "gpt-4o"
DEFAULT_RETRIES = 5
DEFAULT_MIN_SCORE = 70.0
DEFAULT_JUDGE_OUTPUT = "run_outputs/ppt_llm_judge.json"
DEFAULT_COMBINED_REWARD_OUTPUT = "run_outputs/reward_with_ppt_judge.json"
GENERAL_DIMENSIONS = [
    "instruction_following",
    "content_correctness",
    "content_preservation",
    "text_quality",
    "layout_and_readability",
    "visual_design",
    "professional_consistency",
    "no_unnecessary_changes",
]


def _task_path(testbed_dir: str | os.PathLike[str], file_path: str | os.PathLike[str]) -> Path:
    path = Path(file_path)
    if path.is_absolute():
        return path
    return Path(testbed_dir) / path


def _is_valid_pptx(path: Path) -> bool:
    if not path.is_file():
        return False
    try:
        with zipfile.ZipFile(path) as archive:
            return "ppt/presentation.xml" in archive.namelist()
    except zipfile.BadZipFile:
        return False


def summarize_pptx(path: Path) -> dict[str, Any]:
    """Extract a compact text/structure summary from a PPTX file."""
    try:
        from pptx import Presentation
    except ImportError as exc:
        return {
            "file": str(path),
            "error": f"python-pptx is required: {exc}",
            "slides": [],
            "slide_count": None,
        }

    try:
        presentation = Presentation(path)
    except Exception as exc:
        return {
            "file": str(path),
            "error": f"cannot open pptx: {type(exc).__name__}: {exc}",
            "slides": [],
            "slide_count": None,
        }

    slides: list[dict[str, Any]] = []
    for slide_number, slide in enumerate(presentation.slides, start=1):
        slide_text_parts: list[str] = []
        shapes: list[dict[str, Any]] = []
        for shape_index, shape in enumerate(slide.shapes):
            text = ""
            if getattr(shape, "has_text_frame", False):
                text = shape.text or ""
                if text.strip():
                    slide_text_parts.append(text.strip())
            shape_info: dict[str, Any] = {
                "index": shape_index,
                "name": getattr(shape, "name", ""),
                "shape_type": str(getattr(shape, "shape_type", "")),
                "left_pt": _emu_to_pt(getattr(shape, "left", None)),
                "top_pt": _emu_to_pt(getattr(shape, "top", None)),
                "width_pt": _emu_to_pt(getattr(shape, "width", None)),
                "height_pt": _emu_to_pt(getattr(shape, "height", None)),
                "text": text[:1000],
            }
            if getattr(shape, "has_text_frame", False):
                shape_info["runs"] = _summarize_text_runs(shape)
            shapes.append(shape_info)
        slides.append(
            {
                "slide_number": slide_number,
                "text": "\n".join(slide_text_parts)[:3000],
                "shape_count": len(slide.shapes),
                "shapes": shapes[:80],
            }
        )
    return {
        "file": str(path),
        "slide_count": len(presentation.slides),
        "slides": slides,
    }


def _emu_to_pt(value: Any) -> float | None:
    if value is None:
        return None
    try:
        return round(float(value) / 12700.0, 2)
    except (TypeError, ValueError):
        return None


def _summarize_text_runs(shape: Any) -> list[dict[str, Any]]:
    runs: list[dict[str, Any]] = []
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            font = run.font
            runs.append(
                {
                    "text": run.text[:200],
                    "font_name": font.name,
                    "font_size_pt": _font_size_to_pt(font.size),
                    "bold": font.bold,
                    "italic": font.italic,
                    "color_rgb": _font_color(font),
                }
            )
            if len(runs) >= 20:
                return runs
    return runs


def _font_size_to_pt(value: Any) -> float | None:
    if value is None:
        return None
    try:
        return round(float(value.pt), 2)
    except (AttributeError, TypeError, ValueError):
        return None


def _font_color(font: Any) -> str | None:
    try:
        color = font.color.rgb
    except Exception:
        return None
    if color is None:
        return None
    return str(color)[-6:]


def render_pptx_slides(path: Path, max_slides: int = 8) -> dict[str, Any]:
    """Render PPTX slides to PNG images via soffice and pdftoppm when available."""
    soffice = shutil.which("soffice")
    pdftoppm = shutil.which("pdftoppm")
    if not soffice:
        return {"status": "skipped", "reason": "soffice not found", "images": []}
    if not pdftoppm:
        return {"status": "skipped", "reason": "pdftoppm not found", "images": []}

    with tempfile.TemporaryDirectory(prefix="dumate_ppt_judge_") as tmp_name:
        tmp = Path(tmp_name)
        try:
            subprocess.run(
                [soffice, "--headless", "--convert-to", "pdf", "--outdir", str(tmp), str(path)],
                check=True,
                stdout=subprocess.PIPE,
                stderr=subprocess.PIPE,
                text=True,
                timeout=90,
            )
            pdf_path = tmp / f"{path.stem}.pdf"
            if not pdf_path.is_file():
                pdf_candidates = list(tmp.glob("*.pdf"))
                if not pdf_candidates:
                    return {"status": "failed", "reason": "soffice did not produce a PDF", "images": []}
                pdf_path = pdf_candidates[0]
            prefix = tmp / "slide"
            subprocess.run(
                [pdftoppm, "-png", "-r", "120", str(pdf_path), str(prefix)],
                check=True,
                stdout=subprocess.PIPE,
                stderr=subprocess.PIPE,
                text=True,
                timeout=90,
            )
            images = []
            for image_path in sorted(tmp.glob("slide-*.png"))[:max_slides]:
                images.append(
                    {
                        "name": image_path.name,
                        "mime_type": "image/png",
                        "base64": base64.b64encode(image_path.read_bytes()).decode("ascii"),
                    }
                )
            if not images:
                return {"status": "failed", "reason": "pdftoppm did not produce PNG files", "images": []}
            return {"status": "ok", "reason": "", "images": images}
        except subprocess.TimeoutExpired:
            return {"status": "failed", "reason": "rendering timed out", "images": []}
        except subprocess.CalledProcessError as exc:
            detail = (exc.stderr or exc.stdout or str(exc)).strip()
            return {"status": "failed", "reason": detail[:1000], "images": []}


def prepare_evidence(
    testbed_dir: str | os.PathLike[str],
    instruction_file: str,
    output_file: str,
    input_file: str | None = None,
    reference_dir: str | os.PathLike[str] | None = None,
    render_slides: bool = True,
    max_rendered_slides: int = 8,
    reference_summary: dict[str, Any] | None = None,
) -> dict[str, Any]:
    instruction_path = _task_path(testbed_dir, instruction_file)
    output_path = _task_path(testbed_dir, output_file)
    input_path = _task_path(testbed_dir, input_file) if input_file else None
    reference_path = None if reference_dir in {None, "-"} else _task_path(testbed_dir, reference_dir)
    instruction = instruction_path.read_text(encoding="utf-8", errors="ignore") if instruction_path.is_file() else ""
    selected_reference_input: dict[str, Any] | None = None
    if input_path is None and reference_path is not None:
        input_path, selected_reference_input = select_reference_input_ppt(reference_path, instruction, output_path=output_path)

    input_render_status = (
        {"status": "skipped", "reason": "render_slides is false", "images": []}
        if not render_slides
        else {"status": "skipped", "reason": "no input PPT was provided or found in reference_dir", "images": []}
    )
    output_render_status = (
        {"status": "skipped", "reason": "render_slides is false", "images": []}
        if not render_slides
        else {"status": "skipped", "reason": "output PPT is missing or unreadable", "images": []}
    )
    evidence: dict[str, Any] = {
        "instruction_file": str(instruction_path),
        "instruction": instruction,
        "input_file": str(input_path) if input_path else None,
        "output_file": str(output_path),
        "input_valid": bool(input_path and _is_valid_pptx(input_path)),
        "output_valid": _is_valid_pptx(output_path),
        "input_summary": summarize_pptx(input_path) if input_path and input_path.is_file() else None,
        "output_summary": summarize_pptx(output_path) if output_path.is_file() else None,
        "reference_dir": str(reference_path) if reference_path else None,
        "selected_reference_input": selected_reference_input,
        "reference_ppt_summaries": summarize_reference_ppts(reference_path),
        "workspace_reference_summary": reference_summary,
        "render_status": {"input": input_render_status, "output": output_render_status},
    }
    if render_slides:
        if input_path and input_path.is_file():
            evidence["render_status"]["input"] = render_pptx_slides(input_path, max_slides=max_rendered_slides)
        if output_path.is_file():
            evidence["render_status"]["output"] = render_pptx_slides(output_path, max_slides=max_rendered_slides)
    return evidence


def build_ppt_judge_messages(evidence: dict[str, Any], locked_rubrics: list[dict[str, Any]] | None = None) -> list[dict[str, Any]]:
    """Build OpenAI-compatible chat messages for PPT judging."""
    system = (
        "You are a strict evaluator for PowerPoint agent outputs. Judge only from the task "
        "instruction and the two explicitly separated evidence scopes below. Candidate evidence is only the file "
        "under outputs/; reference evidence is only the files under references/. Never treat reference evidence as "
        "candidate output, and never treat candidate evidence as ground truth. If locked task rubrics are provided, "
        "score only those criteria; otherwise first derive task-specific atomic criteria from the general dimensions. "
        "Then score every criterion on an integer 0-4 scale. "
        "Award partial credit, penalize unnecessary changes, and mark critical failures when the output misses "
        "a non-negotiable requirement. Return valid JSON only."
    )
    text_payload = {
        "task_instruction": evidence.get("instruction", ""),
        "general_dimensions": GENERAL_DIMENSIONS,
        "expected_json_schema": {
            "score": "number from 0 to 100",
            "pass": "boolean",
            "task_rubrics": [
                {
                    "id": "snake_case_atomic_id",
                    "dimension": "one general dimension",
                    "description": "one observable task-specific requirement",
                    "weight": "positive number; normalized by the parser",
                    "evidence_required": "boolean",
                    "levels": {
                        "0": "not satisfied or opposite of requirement",
                        "1": "severely deficient",
                        "2": "partially satisfied",
                        "3": "mostly satisfied with minor gaps",
                        "4": "fully satisfied",
                    },
                }
            ],
            "criteria_results": [
                {
                    "id": "same as task_rubrics[].id",
                    "score": "integer 0 to 4",
                    "evidence": "specific slide-level evidence",
                    "rationale": "short reason",
                    "confidence": "number from 0 to 1",
                }
            ],
            "critical_failures": ["list of blocking failures"],
            "summary": "short explanation",
        },
        "candidate_evidence": {
            "scope": "outputs/",
            "output_ppt_summary": _tag_ppt_evidence_paths(evidence.get("output_summary"), "outputs"),
            "render_status": _render_status_without_images(evidence, "output"),
        },
        "reference_evidence": {
            "scope": "references/",
            "workspace_reference_summary": _tag_ppt_evidence_paths(
                evidence.get("workspace_reference_summary"), "references"
            ),
            "reference_ppt_summaries": _tag_ppt_evidence_paths(
                evidence.get("reference_ppt_summaries"), "references"
            ),
            "selected_input_ppt_summary": _tag_ppt_evidence_paths(
                evidence.get("input_summary"), "references"
            ),
            "render_status": _render_status_without_images(evidence, "input"),
        },
        "locked_task_rubrics": locked_rubrics or None,
        "scoring_rules": [
            "If locked_task_rubrics is present, use it exactly as the task_rubrics and do not add, remove, merge, or rewrite criteria.",
            "If locked_task_rubrics is absent, generate 3-16 task_rubrics, each evaluating one observable requirement.",
            "Each task_rubrics item must include weight, evidence_required, and levels 0/1/2/3/4 with task-specific descriptions.",
            "For evidence_required=true criteria, look for ground-truth evidence only in reference_evidence; state what was found or missing in evidence/rationale/confidence.",
            "Every cited candidate path must start with outputs/ and every cited reference path must start with references/.",
            "Each criteria_results score must be an integer from 0 to 4 and reference exactly one rubric id.",
            "The final score must be a weighted 0-100 score computed from criterion scores.",
            "Set pass to false if there are critical failures.",
            "Critical failures include missing/unreadable output, violating explicit slide-count constraints, or failing to preserve core content in an edit task.",
        ],
    }
    content: list[dict[str, Any]] = [{"type": "text", "text": json.dumps(text_payload, ensure_ascii=False, indent=2)}]
    content.extend(_image_content_items(evidence, "input", "references/ selected input PPT rendered slide"))
    content.extend(_image_content_items(evidence, "output", "outputs/ candidate PPT rendered slide"))
    return [{"role": "system", "content": system}, {"role": "user", "content": content}]


def _tag_ppt_evidence_paths(value: Any, scope: str) -> Any:
    """Copy prompt evidence and label file paths with their evidence scope."""
    if isinstance(value, list):
        return [_tag_ppt_evidence_paths(item, scope) for item in value]
    if not isinstance(value, dict):
        return value
    tagged: dict[str, Any] = {}
    for key, item in value.items():
        if key in {"file", "path", "relative_path", "reference_path"} and isinstance(item, str) and item:
            if item.startswith(("outputs/", "references/")):
                tagged[key] = item
            else:
                path = Path(item)
                display = path.name if path.is_absolute() else item.lstrip("/")
                tagged[key] = f"{scope}/{display}"
        else:
            tagged[key] = _tag_ppt_evidence_paths(item, scope)
    return tagged


def _render_status_without_images(evidence: dict[str, Any], key: str) -> dict[str, Any]:
    status = dict(evidence.get("render_status", {}).get(key, {}))
    status["image_count"] = len(status.get("images", []))
    status.pop("images", None)
    return status


def _image_content_items(evidence: dict[str, Any], key: str, label: str) -> list[dict[str, Any]]:
    items: list[dict[str, Any]] = []
    images = evidence.get("render_status", {}).get(key, {}).get("images", [])
    for index, image in enumerate(images, start=1):
        items.append({"type": "text", "text": f"{label} {index}: {image.get('name', '')}"})
        items.append(
            {
                "type": "image_url",
                "image_url": {
                    "url": f"data:{image['mime_type']};base64,{image['base64']}",
                    "detail": "low",
                },
            }
        )
    return items


def aggregate_dimension_score(dimensions: list[dict[str, Any]]) -> float:
    total_weight = 0.0
    earned = 0.0
    for dimension in dimensions:
        try:
            weight = float(dimension.get("weight", 0))
            score = float(dimension.get("score", 0))
        except (TypeError, ValueError):
            continue
        total_weight += weight
        earned += max(0.0, min(5.0, score)) * weight
    if total_weight <= 0:
        return 0.0
    return round((earned / (5.0 * total_weight)) * 100.0, 2)


def _normalize_criteria(rubrics: Any) -> list[dict[str, Any]]:
    if not isinstance(rubrics, list):
        return []
    criteria: list[dict[str, Any]] = []
    seen: set[str] = set()
    for index, item in enumerate(rubrics, start=1):
        if not isinstance(item, dict):
            continue
        cid = str(item.get("id") or f"criterion_{index}").strip()
        if not cid or cid in seen:
            cid = f"criterion_{index}"
        seen.add(cid)
        try:
            weight = float(item.get("weight", 1.0))
        except (TypeError, ValueError):
            weight = 1.0
        levels_raw = item.get("levels") if isinstance(item.get("levels"), dict) else {}
        criteria.append(
            {
                "id": cid,
                "dimension": str(item.get("dimension") or item.get("name") or "instruction_following"),
                "description": str(item.get("description") or item.get("criterion") or item.get("rubric") or cid),
                "weight": max(0.0, weight),
                "evidence_required": bool(item.get("evidence_required", True)),
                "levels": {str(score): str(levels_raw.get(str(score), levels_raw.get(score, ""))) for score in range(5)},
            }
        )
    total = sum(item["weight"] for item in criteria)
    if total <= 0 and criteria:
        total = float(len(criteria))
        for item in criteria:
            item["weight"] = 1.0
    for item in criteria:
        item["weight"] = round(item["weight"] / total, 8) if total else 0.0
    return criteria


def _legacy_dimensions_to_criteria(dimensions: list[dict[str, Any]]) -> tuple[list[dict[str, Any]], list[dict[str, Any]]]:
    criteria = []
    results = []
    for index, dimension in enumerate(dimensions, start=1):
        if not isinstance(dimension, dict):
            continue
        name = str(dimension.get("name") or f"dimension_{index}").strip()
        cid = name or f"dimension_{index}"
        rubric_items = dimension.get("rubric_items")
        description = str(dimension.get("rubric") or dimension.get("evidence") or cid)
        if isinstance(rubric_items, list) and rubric_items:
            first = rubric_items[0] if isinstance(rubric_items[0], dict) else {}
            description = str(first.get("criterion") or description)
        try:
            raw_score = float(dimension.get("score", 0.0))
        except (TypeError, ValueError):
            raw_score = 0.0
        score = int(round(max(0.0, min(5.0, raw_score)) / 5.0 * 4.0))
        criteria.append(
            {
                "id": cid,
                "dimension": name,
                "description": description,
                "weight": float(dimension.get("weight", 1.0) or 1.0),
                "evidence_required": True,
                "levels": {
                    "0": "Not satisfied.",
                    "1": "Severely deficient.",
                    "2": "Partially satisfied.",
                    "3": "Mostly satisfied with minor gaps.",
                    "4": "Fully satisfied.",
                },
            }
        )
        results.append(
            {
                "id": cid,
                "score": score,
                "evidence": str(dimension.get("evidence", "")),
                "rationale": str(dimension.get("rubric", "")),
                "confidence": 0.5,
            }
        )
    return _normalize_criteria(criteria), results


def aggregate_criterion_score(criteria: list[dict[str, Any]], results: list[dict[str, Any]]) -> float:
    result_by_id = {str(item.get("id")): item for item in results if isinstance(item, dict)}
    total_weight = 0.0
    earned = 0.0
    for criterion in criteria:
        cid = str(criterion.get("id", ""))
        try:
            weight = float(criterion.get("weight", 0.0))
            score = int(result_by_id.get(cid, {}).get("score", 0))
        except (TypeError, ValueError):
            continue
        total_weight += max(0.0, weight)
        earned += max(0, min(4, score)) / 4.0 * max(0.0, weight)
    if total_weight <= 0:
        return 0.0
    return round(earned / total_weight * 100.0, 2)


def _parse_json_object_response(content: str) -> dict[str, Any]:
    """Parse a JSON object, tolerating provider-added prose or Markdown fences."""
    try:
        parsed = json.loads(content)
    except json.JSONDecodeError as direct_error:
        decoder = json.JSONDecoder()
        for index, character in enumerate(content):
            if character != "{":
                continue
            try:
                parsed, _ = decoder.raw_decode(content[index:])
            except json.JSONDecodeError:
                continue
            if isinstance(parsed, dict):
                return parsed
        preview = content[:300].replace("\n", "\\n")
        raise ValueError(f"judge response is not valid JSON; response starts with: {preview!r}") from direct_error
    if not isinstance(parsed, dict):
        raise ValueError("judge response must be a JSON object")
    return parsed


def parse_judge_response(raw: str, min_score: float = DEFAULT_MIN_SCORE, locked_rubrics: list[dict[str, Any]] | None = None) -> dict[str, Any]:
    data = _parse_json_object_response(raw.strip())
    dimensions = data.get("dimensions") or []
    criteria = _normalize_criteria(locked_rubrics if locked_rubrics is not None else data.get("task_rubrics"))
    criteria_results = data.get("criteria_results")
    legacy_dimensions = False
    if not criteria and isinstance(dimensions, list) and dimensions:
        legacy_dimensions = True
        criteria, criteria_results = _legacy_dimensions_to_criteria(dimensions)
    if criteria and not isinstance(criteria_results, list):
        criteria_results = []
    if criteria:
        by_id = {item["id"]: item for item in criteria}
        normalized_results = []
        seen: set[str] = set()
        for item in criteria_results if isinstance(criteria_results, list) else []:
            if not isinstance(item, dict):
                continue
            cid = str(item.get("id") or item.get("criterion_id") or "").strip()
            if cid not in by_id or cid in seen:
                continue
            seen.add(cid)
            try:
                score = int(item.get("score"))
            except (TypeError, ValueError):
                score = 0
            score = max(0, min(4, score))
            try:
                confidence = float(item.get("confidence", 0.0))
            except (TypeError, ValueError):
                confidence = 0.0
            normalized_results.append(
                {
                    "id": cid,
                    "score": score,
                    "evidence": str(item.get("evidence", "")),
                    "rationale": str(item.get("rationale", "")),
                    "confidence": max(0.0, min(1.0, confidence)),
                }
            )
        for cid in by_id.keys() - seen:
            normalized_results.append({"id": cid, "score": 0, "evidence": "", "rationale": "Judge did not return this criterion.", "confidence": 0.0})
        data["task_rubrics"] = criteria
        data["criteria_results"] = normalized_results
    if "score" not in data:
        data["score"] = (
            aggregate_dimension_score(dimensions)
            if legacy_dimensions
            else aggregate_criterion_score(criteria, data.get("criteria_results", [])) if criteria else aggregate_dimension_score(dimensions)
        )
    try:
        data["score"] = round(float(data["score"]), 2)
    except (TypeError, ValueError) as exc:
        raise ValueError("judge response score must be numeric") from exc
    data["score"] = max(0.0, min(100.0, data["score"]))
    if not isinstance(data.get("critical_failures"), list):
        data["critical_failures"] = []
    data["pass"] = bool(data["score"] >= min_score and not data["critical_failures"])
    if not isinstance(data.get("dimensions"), list):
        data["dimensions"] = []
    if "summary" not in data:
        data["summary"] = ""
    return data


def call_openai_judge(messages: list[dict[str, Any]], model: str) -> str:
    try:
        from openai import OpenAI
    except ImportError as exc:
        raise RuntimeError("openai>=1.x is required to call the LLM judge") from exc

    base_url = os.environ.get("OPENAI_BASE_URL")
    api_key = os.environ.get("OPENAI_API_KEY")
    if not api_key:
        raise RuntimeError("OPENAI_API_KEY is required to call the LLM judge")
    client = OpenAI(api_key=api_key, base_url=base_url, max_retries=0)
    retries = int(os.environ.get("DUMATE_LLM_JUDGE_RETRIES", str(DEFAULT_RETRIES)))
    last_exc: Exception | None = None
    for attempt in range(max(0, retries) + 1):
        try:
            response = client.chat.completions.create(
                model=model,
                messages=messages,
                temperature=0,
                response_format={"type": "json_object"},
            )
            content = response.choices[0].message.content
            if isinstance(content, list):
                content = "".join(str(part.get("text", "")) for part in content if isinstance(part, dict))
            content = str(content or "").strip()
            if not content:
                raise RuntimeError("PPT LLM judge returned empty content")
            parsed = _parse_json_object_response(content)
            return json.dumps(parsed, ensure_ascii=False)
        except Exception as exc:  # provider SDK exceptions vary by version
            last_exc = exc
            if attempt < max(0, retries):
                logger.warning(
                    "PPT LLM judge request failed on attempt %s/%s: %s",
                    attempt + 1,
                    max(0, retries) + 1,
                    exc,
                )
                time.sleep(1.5 * (attempt + 1))
    raise RuntimeError(
        f"PPT LLM judge request failed after {max(0, retries) + 1} attempts: {last_exc}"
    ) from last_exc


def run_pptx_judge(
    testbed_dir: str | os.PathLike[str],
    instruction_file: str = "instruction.md",
    output_file: str = "",
    input_file: str | None = None,
    reference_dir: str | os.PathLike[str] | None = "workspace_seed",
    model: str | None = None,
    min_score: float = DEFAULT_MIN_SCORE,
    judge_output_file: str = DEFAULT_JUDGE_OUTPUT,
    combined_reward_file: str | None = DEFAULT_COMBINED_REWARD_OUTPUT,
    render_slides: bool = True,
    max_rendered_slides: int = 8,
    reference_summary: dict[str, Any] | None = None,
    locked_rubrics: list[dict[str, Any]] | None = None,
    mock_response: str | None = None,
) -> dict[str, Any]:
    if not output_file:
        raise ValueError("output_file is required")
    model = model or os.environ.get("OPENAI_MODEL") or DEFAULT_MODEL
    evidence = prepare_evidence(
        testbed_dir,
        instruction_file=instruction_file,
        input_file=input_file,
        output_file=output_file,
        reference_dir=reference_dir,
        render_slides=render_slides,
        max_rendered_slides=max_rendered_slides,
        reference_summary=reference_summary,
    )
    status = "ok"
    reason = "PPT LLM judge completed."
    if not evidence["output_valid"]:
        reason = "Candidate PPTX is missing or unreadable."
        result = {
            "score": 0.0,
            "pass": False,
            "dimensions": [],
            "critical_failures": ["missing or unreadable output PPTX"],
            "summary": "The output PowerPoint file is missing or invalid.",
        }
    else:
        normalized_locked_rubrics = _normalize_criteria(locked_rubrics) if locked_rubrics is not None else None
        messages = build_ppt_judge_messages(evidence, normalized_locked_rubrics)
        try:
            raw = mock_response if mock_response is not None else call_openai_judge(messages, model)
            result = parse_judge_response(raw, min_score=min_score, locked_rubrics=normalized_locked_rubrics)
        except Exception as exc:
            logger.warning("PPT LLM judge failed: %s", exc)
            status = "failed"
            reason = f"PPT LLM judge failed: {type(exc).__name__}: {exc}"
            result = {
                "score": None,
                "pass": False,
                "dimensions": [],
                "critical_failures": [],
                "summary": "The PPT LLM judge did not complete; no quality score was assigned.",
            }
    report = {
        "judge_type": "pptx_llm_judge",
        "model": model,
        "status": status,
        "reason": reason,
        "min_score": min_score,
        "result": result,
        "evidence": _evidence_for_report(evidence),
    }
    output_path = _task_path(testbed_dir, judge_output_file)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    output_path.write_text(json.dumps(report, ensure_ascii=False, indent=2), encoding="utf-8")
    if combined_reward_file:
        write_combined_reward(testbed_dir, report, combined_reward_file)
    return report


def _unit_score(value: Any) -> float:
    try:
        score = float(value)
    except (TypeError, ValueError):
        return 0.0
    if score > 1.0:
        score /= 100.0
    return max(0.0, min(1.0, score))


def _read_locked_rubrics(path: str | None) -> list[dict[str, Any]] | None:
    if not path:
        return None
    payload = json.loads(Path(path).expanduser().read_text(encoding="utf-8"))
    if isinstance(payload, dict):
        payload = payload.get("criteria") or payload.get("task_rubrics")
    if not isinstance(payload, list):
        raise ValueError("rubric file must contain a JSON array or an object with criteria/task_rubrics")
    return [item for item in payload if isinstance(item, dict)]


def write_combined_reward(
    testbed_dir: str | os.PathLike[str],
    judge_report: dict[str, Any],
    combined_reward_file: str = DEFAULT_COMBINED_REWARD_OUTPUT,
) -> dict[str, Any] | None:
    reward_path = _task_path(testbed_dir, "run_outputs/reward.json")
    if not reward_path.is_file():
        return None
    try:
        reward = json.loads(reward_path.read_text(encoding="utf-8"))
    except json.JSONDecodeError:
        logger.warning("Cannot combine PPT judge score because reward.json is invalid JSON")
        return None
    if not isinstance(reward, dict):
        return None

    judge_result = judge_report.get("result", {})
    checklist_score = equal_weight_partial_pass(
        reward.get("checks"), fallback=reward.get("partial_pass", 0.0)
    )
    raw_judge_score = judge_result.get("score")
    judge_score = (
        _unit_score(raw_judge_score)
        if judge_report.get("status", "ok") == "ok" and raw_judge_score is not None
        else None
    )
    final_score = (
        merge_final_score(reward.get("complete_pass", 0), checklist_score, judge_score)
        if judge_score is not None
        else None
    )
    min_score = _unit_score(judge_report.get("min_score", DEFAULT_MIN_SCORE))

    combined = dict(reward)
    combined["ppt_llm_judge"] = judge_result
    combined["base_complete_pass"] = reward.get("complete_pass", 0)
    combined["base_partial_pass"] = checklist_score
    combined["ppt_llm_judge_score"] = judge_score
    combined["final_score"] = final_score
    combined["complete_pass_with_ppt_judge"] = int(final_score is not None and final_score >= min_score)
    combined["partial_pass_with_ppt_judge"] = final_score

    combined_path = _task_path(testbed_dir, combined_reward_file)
    combined_path.parent.mkdir(parents=True, exist_ok=True)
    combined_path.write_text(json.dumps(combined, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
    return combined


def _evidence_for_report(evidence: dict[str, Any]) -> dict[str, Any]:
    cleaned = dict(evidence)
    cleaned["render_status"] = {
        "input": _render_status_without_images(evidence, "input"),
        "output": _render_status_without_images(evidence, "output"),
    }
    return cleaned


def summarize_reference_ppts(reference_dir: Path | None, *, max_files: int = 8) -> list[dict[str, Any]]:
    if reference_dir is None or not reference_dir.is_dir():
        return []
    summaries: list[dict[str, Any]] = []
    for path in sorted(reference_dir.rglob("*")):
        if not path.is_file() or path.suffix.lower() not in {".ppt", ".pptx"}:
            continue
        summary = summarize_pptx(path)
        try:
            summary["reference_path"] = path.relative_to(reference_dir).as_posix()
        except ValueError:
            summary["reference_path"] = str(path)
        summaries.append(summary)
        if len(summaries) >= max_files:
            break
    return summaries


def select_reference_input_ppt(
    reference_dir: Path,
    instruction: str,
    *,
    output_path: Path | None = None,
) -> tuple[Path | None, dict[str, Any] | None]:
    """Choose the PPT input most likely required by the task instruction."""
    candidates = _reference_ppt_paths(reference_dir)
    if not candidates:
        return None, None

    instruction_norm = instruction.lower()
    scored: list[tuple[int, str, Path, list[str]]] = []
    output_stem = _normalise_output_stem(output_path.stem) if output_path else ""
    for path in candidates:
        rel = path.relative_to(reference_dir).as_posix()
        name = path.name
        stem = path.stem
        score = 0
        reasons: list[str] = []
        for needle, points, reason in (
            (rel, 100, "relative path appears in instruction"),
            (name, 90, "file name appears in instruction"),
            (stem, 50, "file stem appears in instruction"),
        ):
            if needle and needle.lower() in instruction_norm:
                score += points
                reasons.append(reason)
        if output_stem and _normalise_output_stem(stem) == output_stem:
            score += 20
            reasons.append("file stem matches normalized output stem")
        if not reasons:
            reasons.append("fallback lexical order")
        scored.append((score, rel, path, reasons))

    scored.sort(key=lambda item: (-item[0], item[1]))
    score, rel, path, reasons = scored[0]
    return path, {"path": rel, "score": score, "reasons": reasons}


def _reference_ppt_paths(reference_dir: Path) -> list[Path]:
    if not reference_dir.is_dir():
        return []
    return sorted(
        path for path in reference_dir.rglob("*")
        if path.is_file() and path.suffix.lower() in {".ppt", ".pptx"}
    )


def _normalise_output_stem(stem: str) -> str:
    suffixes = (
        "_优化",
        "-优化",
        " 优化",
        "_edited",
        "-edited",
        "_output",
        "-output",
        "_final",
        "-final",
    )
    value = stem
    changed = True
    while changed:
        changed = False
        lower = value.lower()
        for suffix in suffixes:
            if lower.endswith(suffix.lower()):
                value = value[: -len(suffix)]
                changed = True
                break
    return value.lower()


def evaluate_pptx_llm_judge(testbed_dir: str | os.PathLike[str], args: dict[str, Any]) -> bool:
    try:
        report = run_pptx_judge(
            testbed_dir,
            instruction_file=args.get("instruction_file", "instruction.md"),
            input_file=args.get("input_file"),
            output_file=args["output_file"],
            reference_dir=args.get("reference_dir", "workspace_seed"),
            model=args.get("model"),
            min_score=float(args.get("min_score", DEFAULT_MIN_SCORE)),
            judge_output_file=args.get("judge_output_file", DEFAULT_JUDGE_OUTPUT),
            combined_reward_file=args.get("combined_reward_file", DEFAULT_COMBINED_REWARD_OUTPUT),
            render_slides=bool(args.get("render_slides", True)),
            max_rendered_slides=int(args.get("max_rendered_slides", 8)),
            locked_rubrics=_read_locked_rubrics(args.get("rubric_file")) if args.get("rubric_file") else args.get("locked_rubrics"),
            mock_response=args.get("mock_response"),
        )
    except Exception as exc:
        logger.warning("PPT LLM judge failed: %s", exc)
        return False
    result = report.get("result", {})
    return bool(result.get("pass"))


def main() -> int:
    parser = argparse.ArgumentParser(description="Run DuMateBench PPT LLM-as-judge.")
    parser.add_argument("--task-dir", default=".")
    parser.add_argument("--instruction-file", default="instruction.md")
    parser.add_argument("--input-file")
    parser.add_argument("--output-file", required=True)
    parser.add_argument("--reference-dir", default="workspace_seed", help="Directory containing task input/reference files; use '-' to disable.")
    parser.add_argument("--model", default=None)
    parser.add_argument("--min-score", type=float, default=DEFAULT_MIN_SCORE)
    parser.add_argument("--judge-output-file", default=DEFAULT_JUDGE_OUTPUT)
    parser.add_argument("--combined-reward-file", default=DEFAULT_COMBINED_REWARD_OUTPUT)
    parser.add_argument("--no-render-slides", action="store_true")
    parser.add_argument("--max-rendered-slides", type=int, default=8)
    parser.add_argument("--rubric-file", default="", help="Optional JSON file with fixed criteria/task_rubrics to score instead of generating rubrics.")
    parser.add_argument("--mock-response", help="JSON response for offline smoke tests.")
    args = parser.parse_args()

    report = run_pptx_judge(
        args.task_dir,
        instruction_file=args.instruction_file,
        input_file=args.input_file,
        output_file=args.output_file,
        reference_dir=None if args.reference_dir == "-" else args.reference_dir,
        model=args.model,
        min_score=args.min_score,
        judge_output_file=args.judge_output_file,
        combined_reward_file=args.combined_reward_file,
        render_slides=not args.no_render_slides,
        max_rendered_slides=args.max_rendered_slides,
        locked_rubrics=_read_locked_rubrics(args.rubric_file),
        mock_response=args.mock_response,
    )
    print(json.dumps(report, ensure_ascii=False, indent=2))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
