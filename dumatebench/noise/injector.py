"""Generate deterministic distractor files for DuMateBench workspaces.

The injector intentionally keeps generated noise beside, but separate from, the
source files. It can be used during dataset construction to populate a noisy
workspace while preserving a manifest that identifies every distractor.
"""

from __future__ import annotations

import csv
import json
import math
import random
import struct
import wave
import zipfile
from dataclasses import asdict, dataclass
from pathlib import Path
from typing import Any, Iterable
from xml.sax.saxutils import escape as xml_escape


TEXT_SUFFIXES = {".txt", ".md", ".csv", ".json", ".yaml", ".yml", ".log", ".eml", ".ics"}
EXCEL_SUFFIXES = {".xlsx"}
WORD_SUFFIXES = {".docx"}
PPT_SUFFIXES = {".pptx"}
PDF_SUFFIXES = {".pdf"}
IMAGE_SUFFIXES = {".png", ".jpg", ".jpeg", ".bmp", ".ppm"}
AUDIO_SUFFIXES = {".wav"}
VIDEO_SUFFIXES = {".mp4", ".mov", ".avi", ".mkv"}


@dataclass(frozen=True)
class NoiseConfig:
    """Configuration for deterministic noise generation."""

    output_dir: Path | None = None
    seed: int = 0
    file_noise_count: int = 3
    data_noise_count: int = 2
    include_file_noise: bool = True
    include_data_noise: bool = True
    overwrite: bool = False


@dataclass(frozen=True)
class NoiseRecord:
    """One generated distractor file."""

    source_file: str
    noise_file: str
    category: str
    noise_type: str
    description: str


class NoiseInjector:
    """Create file-level and data-level noise from source files."""

    def __init__(self, config: NoiseConfig | None = None):
        self.config = config or NoiseConfig()
        self._rng = random.Random(self.config.seed)

    def generate(self, files: Iterable[str | Path]) -> dict[str, Any]:
        source_files = [Path(file_path) for file_path in files]
        records: list[NoiseRecord] = []
        for source_file in source_files:
            if not source_file.is_file():
                raise FileNotFoundError(source_file)
            if self.config.include_file_noise:
                records.extend(self._generate_file_noise(source_file))
            if self.config.include_data_noise:
                records.extend(self._generate_data_noise(source_file))

        return {
            "schema_version": "0.1",
            "seed": self.config.seed,
            "target_files": [str(path) for path in source_files],
            "records": [asdict(record) for record in records],
            "distractor_files": [record.noise_file for record in records],
            "noise_types": sorted({record.noise_type for record in records}),
        }

    def _generate_file_noise(self, source_file: Path) -> list[NoiseRecord]:
        patterns = [
            ("historical_version", "{stem}_old{suffix}", "历史版本文件，保留相似命名但内容改为过期项目。"),
            ("temporary_file", "~${stem}.tmp", "临时文件，包含未完成且冲突的数据片段。"),
            ("backup_file", "{stem}.backup{suffix}", "备份文件，模拟自动备份但关键内容明显不同。"),
            ("unrelated_project", "project_notes_{token}.txt", "无关项目文件，包含另一项目的计划和数值。"),
        ]
        records: list[NoiseRecord] = []
        for noise_type, template, description in patterns[: max(0, self.config.file_noise_count)]:
            output_path = self._output_path(source_file, template, force_suffix=source_file.suffix)
            self._write_generic_noise(output_path, source_file, noise_type)
            records.append(self._record(source_file, output_path, "file_noise", noise_type, description))
        return records

    def _generate_data_noise(self, source_file: Path) -> list[NoiseRecord]:
        suffix = source_file.suffix.lower()
        records: list[NoiseRecord] = []
        variants = self._data_variants_for_suffix(suffix)
        for index, noise_type in enumerate(variants[: max(0, self.config.data_noise_count)]):
            output_path = self._output_path(source_file, f"{{stem}}_{noise_type}{{suffix}}")
            writer = self._writer_for_suffix(suffix)
            writer(source_file, output_path, noise_type, index)
            records.append(
                self._record(
                    source_file,
                    output_path,
                    "data_noise",
                    noise_type,
                    self._description_for_data_noise(suffix, noise_type),
                )
            )
        return records

    def _data_variants_for_suffix(self, suffix: str) -> list[str]:
        if suffix in EXCEL_SUFFIXES:
            return ["distractor_sheet", "unit_conversion", "reordered_summary"]
        if suffix in WORD_SUFFIXES or suffix in PPT_SUFFIXES or suffix in PDF_SUFFIXES:
            return ["similar_keywords", "duplicated_content", "truncated_page"]
        if suffix in IMAGE_SUFFIXES:
            return ["visual_noise", "unrelated_material"]
        if suffix in AUDIO_SUFFIXES:
            return ["auditory_noise", "unrelated_material"]
        if suffix in VIDEO_SUFFIXES:
            return ["visual_noise", "unrelated_material"]
        return ["similar_keywords", "duplicated_content", "truncated_content"]

    def _writer_for_suffix(self, suffix: str):
        if suffix in EXCEL_SUFFIXES:
            return self._write_excel_noise
        if suffix in WORD_SUFFIXES:
            return self._write_word_noise
        if suffix in PPT_SUFFIXES:
            return self._write_ppt_noise
        if suffix in PDF_SUFFIXES:
            return self._write_pdf_noise
        if suffix in IMAGE_SUFFIXES:
            return self._write_image_noise
        if suffix in AUDIO_SUFFIXES:
            return self._write_audio_noise
        if suffix in VIDEO_SUFFIXES:
            return self._write_video_noise
        return self._write_text_noise

    def _description_for_data_noise(self, suffix: str, noise_type: str) -> str:
        if suffix in EXCEL_SUFFIXES:
            descriptions = {
                "distractor_sheet": "Excel 干扰 sheet，包含相似表头但不同实体和数值。",
                "unit_conversion": "Excel 单位转换干扰，将数值换算为不同单位并保留相似标签。",
                "reordered_summary": "Excel 重排摘要，改变排序和统计口径。",
            }
            return descriptions[noise_type]
        if suffix in WORD_SUFFIXES or suffix in PPT_SUFFIXES or suffix in PDF_SUFFIXES:
            descriptions = {
                "similar_keywords": "文档干扰内容，包含相似关键字但事实和数值不同。",
                "duplicated_content": "文档重复内容，重复相似段落以干扰检索。",
                "truncated_page": "文档截断页，仅保留不完整上下文。",
            }
            return descriptions[noise_type]
        if suffix in IMAGE_SUFFIXES:
            return "图片干扰素材，包含可见噪点、条纹或无关图像元素。"
        if suffix in AUDIO_SUFFIXES:
            return "音频干扰素材，包含白噪声和提示音。"
        if suffix in VIDEO_SUFFIXES:
            return "视频干扰占位素材，标注为无关视觉材料。"
        return "文本干扰内容，包含相似关键字、重复或截断片段。"

    def _output_path(self, source_file: Path, template: str, *, force_suffix: str | None = None) -> Path:
        base_dir = self.config.output_dir or source_file.parent
        base_dir.mkdir(parents=True, exist_ok=True)
        token = f"{self._rng.randrange(1000, 9999)}"
        suffix = force_suffix if force_suffix is not None and "{suffix}" not in template else source_file.suffix
        name = template.format(stem=source_file.stem, suffix=suffix, token=token)
        path = base_dir / name
        if self.config.overwrite or not path.exists():
            return path
        for index in range(1, 1000):
            candidate = path.with_name(f"{path.stem}_{index}{path.suffix}")
            if not candidate.exists():
                return candidate
        raise FileExistsError(f"Could not choose an unused noise path for {path}")

    def _record(
        self,
        source_file: Path,
        output_path: Path,
        category: str,
        noise_type: str,
        description: str,
    ) -> NoiseRecord:
        return NoiseRecord(
            source_file=str(source_file),
            noise_file=str(output_path),
            category=category,
            noise_type=noise_type,
            description=description,
        )

    def _write_generic_noise(self, output_path: Path, source_file: Path, noise_type: str) -> None:
        suffix = output_path.suffix.lower()
        if suffix in TEXT_SUFFIXES or noise_type == "unrelated_project":
            self._write_text_noise(source_file, output_path, noise_type, 0)
        elif suffix in EXCEL_SUFFIXES:
            self._write_excel_noise(source_file, output_path, "distractor_sheet", 0)
        elif suffix in WORD_SUFFIXES:
            self._write_word_noise(source_file, output_path, "similar_keywords", 0)
        elif suffix in PPT_SUFFIXES:
            self._write_ppt_noise(source_file, output_path, "similar_keywords", 0)
        elif suffix in PDF_SUFFIXES:
            self._write_pdf_noise(source_file, output_path, "similar_keywords", 0)
        elif suffix in IMAGE_SUFFIXES:
            self._write_image_noise(source_file, output_path, "visual_noise", 0)
        elif suffix in AUDIO_SUFFIXES:
            self._write_audio_noise(source_file, output_path, "auditory_noise", 0)
        elif suffix in VIDEO_SUFFIXES:
            self._write_video_noise(source_file, output_path, "visual_noise", 0)
        else:
            output_path.write_bytes(self._generic_binary_payload(source_file, noise_type))

    def _write_text_noise(self, source_file: Path, output_path: Path, noise_type: str, index: int) -> None:
        source_hint = self._source_text_hint(source_file)
        project = self._rng.choice(["Atlas", "Beacon", "Cedar", "Delta", "Evergreen"])
        rows = [
            f"Noise type: {noise_type}",
            f"Reference-looking name: {source_file.name}",
            f"Unrelated project: Project {project}",
            "Status: obsolete draft, do not use for final answers.",
            f"Budget total: {self._rng.randrange(3000, 9000)} USD",
            f"Metric value: {self._rng.randrange(10, 99)}.{self._rng.randrange(10, 99)}",
            f"Similar keywords from source context: {source_hint}",
        ]
        if noise_type in {"duplicated_content", "backup_file"}:
            rows.extend(rows[2:5])
        if noise_type in {"truncated_content", "truncated_page"}:
            rows = rows[:4] + ["[TRUNCATED BEFORE FINAL RESULT]"]

        if output_path.suffix.lower() == ".csv":
            with output_path.open("w", newline="") as handle:
                writer = csv.writer(handle)
                writer.writerow(["field", "value"])
                for line in rows:
                    key, _, value = line.partition(":")
                    writer.writerow([key, value.strip()])
        elif output_path.suffix.lower() == ".json":
            output_path.write_text(json.dumps({"noise": rows}, ensure_ascii=False, indent=2) + "\n")
        else:
            output_path.write_text("\n".join(rows) + "\n")

    def _write_excel_noise(self, source_file: Path, output_path: Path, noise_type: str, index: int) -> None:
        try:
            import openpyxl
        except ImportError:
            rows = self._tabular_noise_rows(source_file, noise_type)
            _write_minimal_xlsx(output_path, self._sheet_name_for_noise(noise_type), rows)
            return

        if source_file.suffix.lower() == ".xlsx":
            try:
                workbook = openpyxl.load_workbook(source_file)
            except Exception:
                workbook = openpyxl.Workbook()
        else:
            workbook = openpyxl.Workbook()
        sheet = workbook.create_sheet(self._unique_sheet_name(workbook, self._sheet_name_for_noise(noise_type)))
        sheet.append(["Item", "Source-like label", "Value", "Unit", "Note"])
        units = [("revenue", "USD", 1), ("revenue", "CNY", 7.2), ("distance", "km", 1.609), ("mass", "lb", 2.204)]
        for row_index in range(1, 8):
            label, unit, multiplier = self._rng.choice(units)
            value = round(self._rng.randrange(50, 500) * multiplier, 2)
            sheet.append(
                [
                    f"Distractor {row_index}",
                    f"{source_file.stem} {label}",
                    value,
                    unit,
                    f"{noise_type}; generated value differs from source",
                ]
            )
        if noise_type == "unit_conversion":
            sheet["G1"] = "Conversion formula"
            sheet["G2"] = "=C2*1.609"
            sheet["H1"] = "Converted unit"
            sheet["H2"] = "alternate"
        workbook.save(output_path)
        workbook.close()

    def _write_word_noise(self, source_file: Path, output_path: Path, noise_type: str, index: int) -> None:
        try:
            from docx import Document
        except ImportError:
            _write_minimal_docx(output_path, [f"{source_file.stem} - distractor draft"] + self._document_noise_paragraphs(source_file, noise_type))
            return

        document = Document()
        document.add_heading(f"{source_file.stem} - distractor draft", level=1)
        for paragraph in self._document_noise_paragraphs(source_file, noise_type):
            document.add_paragraph(paragraph)
        document.save(output_path)

    def _write_ppt_noise(self, source_file: Path, output_path: Path, noise_type: str, index: int) -> None:
        try:
            from pptx import Presentation
            from pptx.util import Inches
        except ImportError:
            _write_minimal_pptx(output_path, f"{source_file.stem} distractor", self._document_noise_paragraphs(source_file, noise_type))
            return

        presentation = Presentation()
        paragraphs = self._document_noise_paragraphs(source_file, noise_type)
        for title, body in [(f"{source_file.stem} distractor", paragraphs[0]), ("Notes", "\n".join(paragraphs[1:]))]:
            slide = presentation.slides.add_slide(presentation.slide_layouts[6])
            title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.5), Inches(8.0), Inches(0.8))
            title_box.text_frame.text = title
            body_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(8.2), Inches(4.5))
            body_box.text_frame.text = body
        presentation.save(output_path)

    def _write_pdf_noise(self, source_file: Path, output_path: Path, noise_type: str, index: int) -> None:
        lines = [f"{source_file.stem} distractor"] + self._document_noise_paragraphs(source_file, noise_type)
        _write_minimal_pdf(output_path, lines[:8])

    def _write_image_noise(self, source_file: Path, output_path: Path, noise_type: str, index: int) -> None:
        if output_path.suffix.lower() == ".bmp":
            _write_noise_bmp(output_path, self._rng, width=96, height=64)
            return
        if output_path.suffix.lower() == ".ppm":
            _write_noise_ppm(output_path, self._rng, width=96, height=64)
            return
        # Keep a valid image even without Pillow by using PPM bytes under the
        # requested name; consumers that sniff content can still read it.
        _write_noise_ppm(output_path, self._rng, width=96, height=64)

    def _write_audio_noise(self, source_file: Path, output_path: Path, noise_type: str, index: int) -> None:
        sample_rate = 8000
        duration_sec = 1.0
        with wave.open(str(output_path), "wb") as wav:
            wav.setnchannels(1)
            wav.setsampwidth(2)
            wav.setframerate(sample_rate)
            frames = bytearray()
            for sample_index in range(int(sample_rate * duration_sec)):
                tone = math.sin(2 * math.pi * 440 * sample_index / sample_rate) * 0.35
                hiss = self._rng.uniform(-0.2, 0.2)
                sample = int(max(-1.0, min(1.0, tone + hiss)) * 32767)
                frames.extend(struct.pack("<h", sample))
            wav.writeframes(bytes(frames))

    def _write_video_noise(self, source_file: Path, output_path: Path, noise_type: str, index: int) -> None:
        output_path.write_bytes(
            b"DUMATEBENCH-UNRELATED-VIDEO-NOISE\n"
            + f"source={source_file.name}\nnoise_type={noise_type}\n".encode()
            + self._rng.randbytes(256)
        )

    def _write_tabular_text_fallback(self, output_path: Path, source_file: Path, noise_type: str) -> None:
        rows = self._tabular_noise_rows(source_file, noise_type)
        with output_path.open("w", newline="") as handle:
            writer = csv.writer(handle)
            writer.writerows(rows)

    def _tabular_noise_rows(self, source_file: Path, noise_type: str) -> list[list[Any]]:
        rows = [
            ["Item", "Source-like label", "Value", "Unit", "Note"],
            ["Distractor 1", source_file.stem, self._rng.randrange(100, 999), "alternate", noise_type],
            ["Distractor 2", f"{source_file.stem} old", self._rng.randrange(100, 999), "converted", "different"],
        ]
        return rows

    def _document_noise_paragraphs(self, source_file: Path, noise_type: str) -> list[str]:
        hint = self._source_text_hint(source_file)
        paragraphs = [
            f"This file intentionally resembles {source_file.name}, but it belongs to a different scenario.",
            f"Similar keywords: {hint}. Final answer values here are obsolete and should not match the source.",
            f"Changed total: {self._rng.randrange(1000, 9000)}; changed date: 2025-{self._rng.randrange(1, 13):02d}-{self._rng.randrange(1, 28):02d}.",
        ]
        if noise_type == "duplicated_content":
            paragraphs.extend(paragraphs[1:])
        if noise_type == "truncated_page":
            paragraphs = paragraphs[:2] + ["[TRUNCATED PAGE: remaining context intentionally missing]"]
        return paragraphs

    def _source_text_hint(self, source_file: Path, max_words: int = 12) -> str:
        try:
            text = source_file.read_text(errors="ignore")
        except Exception:
            text = source_file.stem.replace("_", " ")
        text = _xml_compatible_text(text)
        words = [word.strip(".,:;()[]{}") for word in text.split() if word.strip()]
        if not words:
            words = source_file.stem.replace("_", " ").split()
        return " ".join(words[:max_words]) or source_file.stem

    def _sheet_name_for_noise(self, noise_type: str) -> str:
        return {
            "distractor_sheet": "Archive_Data",
            "unit_conversion": "Unit_Conversion",
            "reordered_summary": "Old_Summary",
        }.get(noise_type, "Noise")

    def _unique_sheet_name(self, workbook: Any, name: str) -> str:
        existing = set(workbook.sheetnames)
        if name not in existing:
            return name
        for index in range(1, 100):
            candidate = f"{name}_{index}"
            if candidate not in existing:
                return candidate
        return f"{name}_{self._rng.randrange(1000, 9999)}"

    def _generic_binary_payload(self, source_file: Path, noise_type: str) -> bytes:
        header = f"DUMATEBENCH NOISE\nsource={source_file.name}\ntype={noise_type}\n".encode()
        return header + self._rng.randbytes(512)


def inject_noise(files: Iterable[str | Path], config: NoiseConfig | None = None) -> dict[str, Any]:
    """Generate noise for files and return a manifest dictionary."""

    return NoiseInjector(config).generate(files)


def _xml_compatible_text(text: str) -> str:
    """Replace characters that cannot be written into XML-backed Office files."""

    cleaned = []
    for char in text:
        codepoint = ord(char)
        if (
            char in {"\t", "\n", "\r"}
            or 0x20 <= codepoint <= 0xD7FF
            or 0xE000 <= codepoint <= 0xFFFD
            or 0x10000 <= codepoint <= 0x10FFFF
        ):
            cleaned.append(char)
        else:
            cleaned.append(" ")
    return "".join(cleaned)


def _write_minimal_pdf(path: Path, lines: list[str]) -> None:
    escaped_lines = [_escape_pdf_text(line[:100]) for line in lines]
    text_ops = ["BT", "/F1 12 Tf", "72 740 Td"]
    for index, line in enumerate(escaped_lines):
        if index:
            text_ops.append("0 -18 Td")
        text_ops.append(f"({line}) Tj")
    text_ops.append("ET")
    stream = "\n".join(text_ops).encode("latin-1", errors="replace")
    objects = [
        b"<< /Type /Catalog /Pages 2 0 R >>",
        b"<< /Type /Pages /Kids [3 0 R] /Count 1 >>",
        b"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] /Resources << /Font << /F1 4 0 R >> >> /Contents 5 0 R >>",
        b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>",
        b"<< /Length " + str(len(stream)).encode() + b" >>\nstream\n" + stream + b"\nendstream",
    ]
    content = bytearray(b"%PDF-1.4\n")
    offsets = [0]
    for number, obj in enumerate(objects, start=1):
        offsets.append(len(content))
        content.extend(f"{number} 0 obj\n".encode())
        content.extend(obj)
        content.extend(b"\nendobj\n")
    xref_offset = len(content)
    content.extend(f"xref\n0 {len(objects) + 1}\n".encode())
    content.extend(b"0000000000 65535 f \n")
    for offset in offsets[1:]:
        content.extend(f"{offset:010d} 00000 n \n".encode())
    content.extend(
        f"trailer\n<< /Size {len(objects) + 1} /Root 1 0 R >>\nstartxref\n{xref_offset}\n%%EOF\n".encode()
    )
    path.write_bytes(bytes(content))


def _escape_pdf_text(text: str) -> str:
    return text.replace("\\", "\\\\").replace("(", "\\(").replace(")", "\\)")


def _write_minimal_xlsx(path: Path, sheet_name: str, rows: list[list[Any]]) -> None:
    sheet_data: list[str] = []
    for row_index, row in enumerate(rows, start=1):
        cells = []
        for col_index, value in enumerate(row, start=1):
            cell_ref = f"{_excel_column(col_index)}{row_index}"
            if isinstance(value, int | float):
                cells.append(f'<c r="{cell_ref}"><v>{value}</v></c>')
            else:
                cells.append(
                    f'<c r="{cell_ref}" t="inlineStr"><is><t>{xml_escape(str(value))}</t></is></c>'
                )
        sheet_data.append(f'<row r="{row_index}">{"".join(cells)}</row>')
    with zipfile.ZipFile(path, "w", compression=zipfile.ZIP_DEFLATED) as archive:
        archive.writestr(
            "[Content_Types].xml",
            """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">
  <Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/>
  <Default Extension="xml" ContentType="application/xml"/>
  <Override PartName="/xl/workbook.xml" ContentType="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet.main+xml"/>
  <Override PartName="/xl/worksheets/sheet1.xml" ContentType="application/vnd.openxmlformats-officedocument.spreadsheetml.worksheet+xml"/>
</Types>""",
        )
        archive.writestr(
            "_rels/.rels",
            """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
  <Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" Target="xl/workbook.xml"/>
</Relationships>""",
        )
        archive.writestr(
            "xl/workbook.xml",
            f"""<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<workbook xmlns="http://schemas.openxmlformats.org/spreadsheetml/2006/main"
          xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
  <sheets><sheet name="{xml_escape(sheet_name[:31])}" sheetId="1" r:id="rId1"/></sheets>
</workbook>""",
        )
        archive.writestr(
            "xl/_rels/workbook.xml.rels",
            """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
  <Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/worksheet" Target="worksheets/sheet1.xml"/>
</Relationships>""",
        )
        archive.writestr(
            "xl/worksheets/sheet1.xml",
            f"""<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<worksheet xmlns="http://schemas.openxmlformats.org/spreadsheetml/2006/main">
  <sheetData>{''.join(sheet_data)}</sheetData>
</worksheet>""",
        )


def _write_minimal_docx(path: Path, paragraphs: list[str]) -> None:
    body = "".join(
        f"<w:p><w:r><w:t>{xml_escape(paragraph)}</w:t></w:r></w:p>" for paragraph in paragraphs
    )
    with zipfile.ZipFile(path, "w", compression=zipfile.ZIP_DEFLATED) as archive:
        archive.writestr(
            "[Content_Types].xml",
            """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">
  <Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/>
  <Default Extension="xml" ContentType="application/xml"/>
  <Override PartName="/word/document.xml" ContentType="application/vnd.openxmlformats-officedocument.wordprocessingml.document.main+xml"/>
</Types>""",
        )
        archive.writestr(
            "_rels/.rels",
            """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
  <Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" Target="word/document.xml"/>
</Relationships>""",
        )
        archive.writestr(
            "word/document.xml",
            f"""<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<w:document xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main">
  <w:body>{body}<w:sectPr/></w:body>
</w:document>""",
        )


def _write_minimal_pptx(path: Path, title: str, paragraphs: list[str]) -> None:
    body_text = " ".join(paragraphs)
    with zipfile.ZipFile(path, "w", compression=zipfile.ZIP_DEFLATED) as archive:
        archive.writestr(
            "[Content_Types].xml",
            """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">
  <Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/>
  <Default Extension="xml" ContentType="application/xml"/>
  <Override PartName="/ppt/presentation.xml" ContentType="application/vnd.openxmlformats-officedocument.presentationml.presentation.main+xml"/>
  <Override PartName="/ppt/slides/slide1.xml" ContentType="application/vnd.openxmlformats-officedocument.presentationml.slide+xml"/>
</Types>""",
        )
        archive.writestr(
            "_rels/.rels",
            """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
  <Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" Target="ppt/presentation.xml"/>
</Relationships>""",
        )
        archive.writestr(
            "ppt/presentation.xml",
            """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<p:presentation xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
                xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
  <p:sldIdLst><p:sldId id="256" r:id="rId1"/></p:sldIdLst>
</p:presentation>""",
        )
        archive.writestr(
            "ppt/_rels/presentation.xml.rels",
            """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
  <Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide" Target="slides/slide1.xml"/>
</Relationships>""",
        )
        archive.writestr(
            "ppt/slides/slide1.xml",
            f"""<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<p:sld xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
       xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
  <p:cSld><p:spTree>
    <p:nvGrpSpPr><p:cNvPr id="1" name=""/><p:cNvGrpSpPr/><p:nvPr/></p:nvGrpSpPr><p:grpSpPr/>
    <p:sp><p:nvSpPr><p:cNvPr id="2" name="Title"/><p:cNvSpPr/><p:nvPr/></p:nvSpPr>
      <p:txBody><a:bodyPr/><a:lstStyle/><a:p><a:r><a:t>{xml_escape(title)}</a:t></a:r></a:p></p:txBody>
    </p:sp>
    <p:sp><p:nvSpPr><p:cNvPr id="3" name="Body"/><p:cNvSpPr/><p:nvPr/></p:nvSpPr>
      <p:txBody><a:bodyPr/><a:lstStyle/><a:p><a:r><a:t>{xml_escape(body_text)}</a:t></a:r></a:p></p:txBody>
    </p:sp>
  </p:spTree></p:cSld>
</p:sld>""",
        )


def _excel_column(index: int) -> str:
    name = ""
    while index:
        index, remainder = divmod(index - 1, 26)
        name = chr(65 + remainder) + name
    return name


def _write_noise_ppm(path: Path, rng: random.Random, *, width: int, height: int) -> None:
    data = bytearray(f"P6\n{width} {height}\n255\n".encode())
    for y in range(height):
        for x in range(width):
            stripe = 220 if (x // 8 + y // 8) % 2 == 0 else 40
            data.extend(
                [
                    (stripe + rng.randrange(0, 36)) % 256,
                    rng.randrange(0, 256),
                    (255 - stripe + rng.randrange(0, 36)) % 256,
                ]
            )
    path.write_bytes(bytes(data))


def _write_noise_bmp(path: Path, rng: random.Random, *, width: int, height: int) -> None:
    row_size = (width * 3 + 3) & ~3
    pixel_data = bytearray()
    for _y in range(height):
        row = bytearray()
        for _x in range(width):
            row.extend([rng.randrange(256), rng.randrange(256), rng.randrange(256)])
        row.extend(b"\x00" * (row_size - width * 3))
        pixel_data.extend(row)
    file_size = 54 + len(pixel_data)
    header = bytearray()
    header.extend(b"BM")
    header.extend(struct.pack("<IHHI", file_size, 0, 0, 54))
    header.extend(struct.pack("<IIIHHIIIIII", 40, width, height, 1, 24, 0, len(pixel_data), 2835, 2835, 0, 0))
    path.write_bytes(bytes(header + pixel_data))
